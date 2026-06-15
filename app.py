import os
import io
import json
import uuid
import smtplib
import threading
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
from datetime import datetime, timedelta
from functools import wraps
from flask import (Flask, render_template, request, redirect, url_for,
                   flash, session, jsonify, send_file)
from flask_sqlalchemy import SQLAlchemy
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
# Fix CLOUDINARY_URL before importing the library — it auto-reads the env var
# at import time and crashes if the value is invalid
_raw_cld = os.environ.get('CLOUDINARY_URL', '')
if _raw_cld and not _raw_cld.startswith('cloudinary://'):
    # User may have pasted the full line "CLOUDINARY_URL=cloudinary://..."
    if '=' in _raw_cld:
        _raw_cld = _raw_cld.split('=', 1)[1].strip()
    if _raw_cld.startswith('cloudinary://'):
        os.environ['CLOUDINARY_URL'] = _raw_cld   # fixed
    else:
        del os.environ['CLOUDINARY_URL']           # remove bad value entirely
import cloudinary
import cloudinary.uploader
import cloudinary.api

app = Flask(__name__)
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'autopow3r-s3cr3t-k3y-2024-change-me')
_db_url = os.environ.get('DATABASE_URL', 'sqlite:///dealership.db')
if _db_url.startswith('postgres://'):          # Render gives postgres://, SQLAlchemy needs postgresql://
    _db_url = _db_url.replace('postgres://', 'postgresql://', 1)
app.config['SQLALCHEMY_DATABASE_URI'] = _db_url
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['UPLOAD_FOLDER'] = os.path.join('static', 'uploads')
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500 MB

BUSINESS = {
    'name': 'AutoPower Dealership',
    'tagline': 'Drive Excellence. Power the Future.',
    'phone': '+233 50 356 6913',
    'email': 'kofi@chinacarsinghana.com',
    'address': 'Accra, Ghana',
    'republic_bank_url': '#',
    'whatsapp': '233503566913',
    'facebook': '#',
    'instagram': '#',
}

ALLOWED_IMAGES = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
ALLOWED_VIDEOS = {'mp4', 'mov', 'avi', 'mkv', 'webm'}
ALLOWED_ALL = ALLOWED_IMAGES | ALLOWED_VIDEOS

# ── Email configuration ───────────────────────────────────────────────────── #
LOAN_TO_EMAIL       = 'sbonsu@republicghana.com'          # Republic Bank primary
LOAN_CC_EMAILS      = ['kofi@chinacarsinghana.com', 'pdesbordes@republicghana.com']
ACCESS_BANK_EMAILS  = ['emmanueloffei44@gmail.com', 'Offeie@accessbankplc.com']
ACCESS_BANK_CC      = ['kofi@chinacarsinghana.com']
NOTIFY_EMAILS       = ['kofi@chinacarsinghana.com']  # orders, test drives, contacts
# SMTP settings
SMTP_HOST     = 'smtp.gmail.com'
SMTP_PORT     = 587
SMTP_USER     = os.environ.get('SMTP_USER', '')
SMTP_PASSWORD = os.environ.get('SMTP_PASSWORD', '')

# Cloudinary — persistent cloud storage for uploaded images/videos
USE_CLOUDINARY = False
try:
    _cld_url = os.environ.get('CLOUDINARY_URL', '').strip()
    # Strip the variable name if user accidentally pasted "CLOUDINARY_URL=cloudinary://..."
    if '=' in _cld_url and _cld_url.startswith('CLOUDINARY_URL'):
        _cld_url = _cld_url.split('=', 1)[1].strip()
    if _cld_url.startswith('cloudinary://'):
        cloudinary.config(cloudinary_url=_cld_url)
        USE_CLOUDINARY = True
    elif os.environ.get('CLOUDINARY_CLOUD_NAME'):
        cloudinary.config(
            cloud_name = os.environ.get('CLOUDINARY_CLOUD_NAME', ''),
            api_key    = os.environ.get('CLOUDINARY_API_KEY', ''),
            api_secret = os.environ.get('CLOUDINARY_API_SECRET', ''),
        )
        USE_CLOUDINARY = True
    else:
        print('[CLOUDINARY] Not configured — using local file storage.')
except Exception as e:
    print(f'[CLOUDINARY] Config error: {e} — falling back to local storage.')


def send_loan_notification(app_obj):
    """Email loan application details to the chosen bank and internal CC."""
    if not SMTP_USER or not SMTP_PASSWORD:
        return
    try:
        rate         = get_ghs_rate()
        ptype        = (app_obj.product_type or '').title()
        pname        = app_obj.product_name or 'Not specified'
        p_usd        = f"${app_obj.product_price:,.2f}"  if app_obj.product_price else 'N/A'
        p_ghs        = f"GH₵ {app_obj.product_price * rate:,.2f}" if app_obj.product_price else 'N/A'
        l_usd        = f"${app_obj.loan_amount:,.2f}"    if app_obj.loan_amount    else 'N/A'
        l_ghs        = f"GH₵ {app_obj.loan_amount * rate:,.2f}"   if app_obj.loan_amount    else 'N/A'
        d_usd        = f"${app_obj.deposit_amount:,.2f}" if app_obj.deposit_amount else 'N/A'
        d_ghs        = f"GH₵ {app_obj.deposit_amount * rate:,.2f}" if app_obj.deposit_amount else 'N/A'
        income_str   = f"${app_obj.annual_income:,.2f} (GH₵ {app_obj.annual_income * rate:,.2f})" if app_obj.annual_income else 'N/A'

        is_access = (app_obj.bank_choice or 'republic') == 'access'
        bank_label = 'ACCESS BANK GHANA' if is_access else 'REPUBLIC BANK GHANA'

        sp_section = ''
        if app_obj.sp_code:
            sp_section = f"""
  REFERRED BY SALESPERSON
  -----------------------
  Name  : {app_obj.sp_name}
  Code  : {app_obj.sp_code}
  Phone : {app_obj.sp_phone or 'N/A'}
"""
        body = f"""\
NEW LOAN APPLICATION — #{app_obj.id}
{'=' * 54}
  BANK SELECTED : {bank_label}

  PRODUCT REQUESTED
  -----------------
  Type   : {ptype}
  Name   : {pname}
  Price  : {p_usd}  |  {p_ghs}

  LOAN DETAILS
  ------------
  Loan Amount : {l_usd}  |  {l_ghs}
  Deposit     : {d_usd}  |  {d_ghs}
  Loan Term   : {app_obj.loan_term or 'N/A'} months

  APPLICANT DETAILS
  -----------------
  Full Name     : {app_obj.first_name} {app_obj.last_name}
  Email         : {app_obj.email}
  Phone         : {app_obj.phone}
  Date of Birth : {app_obj.date_of_birth or 'N/A'}
  ID / Passport : {app_obj.id_number or 'N/A'}
  Address       : {app_obj.address or 'N/A'}

  EMPLOYMENT & INCOME
  -------------------
  Employment  : {app_obj.employment_status or 'N/A'}
  Employer    : {app_obj.employer or 'N/A'}
  Ann. Income : {income_str}
{sp_section}
  ADDITIONAL NOTES
  ----------------
  {app_obj.message or 'None'}

{'=' * 54}
Submitted : {app_obj.created_at.strftime('%d %B %Y at %H:%M UTC')}
"""
        if is_access:
            to_addr    = ACCESS_BANK_EMAILS[0]
            cc_addrs   = ACCESS_BANK_EMAILS[1:] + ACCESS_BANK_CC
        else:
            to_addr    = LOAN_TO_EMAIL
            cc_addrs   = LOAN_CC_EMAILS

        all_recipients = [to_addr] + cc_addrs
        bank_tag = 'Access Bank' if is_access else 'Republic Bank'
        msg = MIMEMultipart('alternative')
        msg['Subject'] = f"[{bank_tag} Loan #{app_obj.id}] {ptype}: {pname} — {app_obj.first_name} {app_obj.last_name}"
        msg['From']    = SMTP_USER
        msg['To']      = to_addr
        msg['Cc']      = ', '.join(cc_addrs)
        msg.attach(MIMEText(body, 'plain'))

        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
            server.ehlo()
            server.starttls()
            server.login(SMTP_USER, SMTP_PASSWORD)
            server.sendmail(SMTP_USER, all_recipients, msg.as_string())
    except Exception as e:
        print(f'[EMAIL] Failed to send loan notification: {e}')


db = SQLAlchemy(app)


# ─────────────────────────────── MODELS ─────────────────────────────────── #

class Vehicle(db.Model):
    id            = db.Column(db.Integer, primary_key=True)
    name          = db.Column(db.String(100), nullable=False)
    tagline       = db.Column(db.String(200))
    description   = db.Column(db.Text)
    price         = db.Column(db.Float, nullable=False, default=0)
    category      = db.Column(db.String(50))   # SUV, Sedan, Truck, Van, Pickup
    year          = db.Column(db.Integer)
    mileage       = db.Column(db.Integer)
    transmission  = db.Column(db.String(50))   # Automatic, Manual
    fuel_type     = db.Column(db.String(50))   # Petrol, Diesel, Electric, Hybrid
    engine        = db.Column(db.String(100))
    color         = db.Column(db.String(50))
    seats         = db.Column(db.Integer)
    features      = db.Column(db.Text, default='[]')
    images        = db.Column(db.Text, default='[]')
    videos        = db.Column(db.Text, default='[]')
    featured      = db.Column(db.Boolean, default=False)
    available     = db.Column(db.Boolean, default=True)
    created_at    = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at    = db.Column(db.DateTime, default=datetime.utcnow)

    def get_images(self):
        try: return json.loads(self.images or '[]')
        except: return []

    def get_videos(self):
        try: return json.loads(self.videos or '[]')
        except: return []

    def get_features(self):
        try: return json.loads(self.features or '[]')
        except: return []

    @property
    def cover(self):
        imgs = self.get_images()
        return imgs[0] if imgs else None


class SolarSystem(db.Model):
    id               = db.Column(db.Integer, primary_key=True)
    name             = db.Column(db.String(100), nullable=False)
    tagline          = db.Column(db.String(200))
    description      = db.Column(db.Text)
    price            = db.Column(db.Float, nullable=False, default=0)
    category         = db.Column(db.String(50))   # Residential, Commercial, Industrial
    power_output     = db.Column(db.String(50))
    panel_count      = db.Column(db.Integer)
    battery_capacity = db.Column(db.String(50))
    warranty         = db.Column(db.String(50))
    efficiency       = db.Column(db.String(50))
    features         = db.Column(db.Text, default='[]')
    images           = db.Column(db.Text, default='[]')
    videos           = db.Column(db.Text, default='[]')
    featured         = db.Column(db.Boolean, default=False)
    available        = db.Column(db.Boolean, default=True)
    created_at       = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at       = db.Column(db.DateTime, default=datetime.utcnow)

    def get_images(self):
        try: return json.loads(self.images or '[]')
        except: return []

    def get_videos(self):
        try: return json.loads(self.videos or '[]')
        except: return []

    def get_features(self):
        try: return json.loads(self.features or '[]')
        except: return []

    @property
    def cover(self):
        imgs = self.get_images()
        return imgs[0] if imgs else None


class LoanApplication(db.Model):
    id                = db.Column(db.Integer, primary_key=True)
    first_name        = db.Column(db.String(50),  nullable=False)
    last_name         = db.Column(db.String(50),  nullable=False)
    email             = db.Column(db.String(100), nullable=False)
    phone             = db.Column(db.String(20),  nullable=False)
    date_of_birth     = db.Column(db.String(20))
    id_number         = db.Column(db.String(50))
    address           = db.Column(db.Text)
    product_type      = db.Column(db.String(20))   # vehicle | solar
    product_id        = db.Column(db.Integer)
    product_name      = db.Column(db.String(100))
    product_price     = db.Column(db.Float)
    loan_amount       = db.Column(db.Float)
    deposit_amount    = db.Column(db.Float)
    loan_term         = db.Column(db.Integer)       # months
    employment_status = db.Column(db.String(50))
    employer          = db.Column(db.String(100))
    annual_income     = db.Column(db.Float)
    message           = db.Column(db.Text)
    bank_choice       = db.Column(db.String(20), default='republic')  # 'republic' | 'access'
    status            = db.Column(db.String(20), default='pending')
    created_at        = db.Column(db.DateTime, default=datetime.utcnow)
    sp_code           = db.Column(db.String(20))
    sp_name           = db.Column(db.String(100))
    sp_phone          = db.Column(db.String(20))


class Contact(db.Model):
    id         = db.Column(db.Integer, primary_key=True)
    name       = db.Column(db.String(100), nullable=False)
    email      = db.Column(db.String(100), nullable=False)
    phone      = db.Column(db.String(20))
    subject    = db.Column(db.String(200))
    message    = db.Column(db.Text, nullable=False)
    status     = db.Column(db.String(20), default='unread')
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    sp_code    = db.Column(db.String(20))
    sp_name    = db.Column(db.String(100))
    sp_phone   = db.Column(db.String(20))
    reply_text = db.Column(db.Text)
    replied_at = db.Column(db.DateTime)


class CarOrder(db.Model):
    id             = db.Column(db.Integer, primary_key=True)
    order_type     = db.Column(db.String(20), nullable=False)  # 'order' | 'test_drive'
    vehicle_id     = db.Column(db.Integer, db.ForeignKey('vehicle.id'), nullable=True)
    vehicle        = db.relationship('Vehicle', backref='orders')
    name           = db.Column(db.String(100), nullable=False)
    email          = db.Column(db.String(100), nullable=False)
    phone          = db.Column(db.String(20),  nullable=False)
    location       = db.Column(db.String(200))   # user's location / address
    latitude       = db.Column(db.Float)
    longitude      = db.Column(db.Float)
    preferred_date = db.Column(db.String(20))
    preferred_time = db.Column(db.String(20))
    notes          = db.Column(db.Text)
    status         = db.Column(db.String(20), default='pending')
    created_at     = db.Column(db.DateTime, default=datetime.utcnow)
    sp_code        = db.Column(db.String(20))
    sp_name        = db.Column(db.String(100))
    sp_phone       = db.Column(db.String(20))


class AdminUser(db.Model):
    id            = db.Column(db.Integer, primary_key=True)
    username      = db.Column(db.String(50), unique=True, nullable=False)
    password_hash = db.Column(db.String(200), nullable=False)
    email         = db.Column(db.String(100))
    full_name     = db.Column(db.String(100))
    created_at    = db.Column(db.DateTime, default=datetime.utcnow)

    def set_password(self, pw):
        self.password_hash = generate_password_hash(pw)

    def check_password(self, pw):
        return check_password_hash(self.password_hash, pw)


class Salesperson(db.Model):
    id           = db.Column(db.Integer, primary_key=True)
    full_name    = db.Column(db.String(100), nullable=False)
    email        = db.Column(db.String(100))
    phone        = db.Column(db.String(20))
    code         = db.Column(db.String(20), unique=True, nullable=False)  # unique referral/access code
    password_hash      = db.Column(db.String(200), nullable=False)
    active             = db.Column(db.Boolean, default=True)
    commission_pct     = db.Column(db.Float, default=1.0)  # commission percentage
    notes              = db.Column(db.Text)
    created_at         = db.Column(db.DateTime, default=datetime.utcnow)
    reset_token        = db.Column(db.String(64))
    reset_token_expiry = db.Column(db.DateTime)
    # stats (computed from linked applications/sales)

    def set_password(self, pw):
        self.password_hash = generate_password_hash(pw)

    def check_password(self, pw):
        return check_password_hash(self.password_hash, pw)

    def total_sales(self):
        return SalespersonSale.query.filter_by(salesperson_id=self.id).filter(
            SalespersonSale.status != 'cancelled').count()

    def total_revenue(self):
        sales = SalespersonSale.query.filter_by(salesperson_id=self.id).filter(
            SalespersonSale.status != 'cancelled').all()
        return sum(s.sale_amount or 0 for s in sales)

    def commission_earned(self):
        return self.total_revenue() * (self.commission_pct / 100)


class SalespersonSale(db.Model):
    id               = db.Column(db.Integer, primary_key=True)
    salesperson_id   = db.Column(db.Integer, db.ForeignKey('salesperson.id'), nullable=False)
    salesperson      = db.relationship('Salesperson', backref='sales')
    product_type     = db.Column(db.String(20))   # vehicle | solar
    product_id       = db.Column(db.Integer)
    product_name     = db.Column(db.String(100))
    sale_amount      = db.Column(db.Float)
    customer_name    = db.Column(db.String(100))
    customer_email   = db.Column(db.String(100))
    customer_phone   = db.Column(db.String(20))
    loan_application_id = db.Column(db.Integer, db.ForeignKey('loan_application.id'), nullable=True)
    notes            = db.Column(db.Text)
    status           = db.Column(db.String(20), default='pending')  # pending, confirmed, cancelled
    created_at       = db.Column(db.DateTime, default=datetime.utcnow)


class SiteSettings(db.Model):
    key   = db.Column(db.String(50), primary_key=True)
    value = db.Column(db.String(200), nullable=False)


class SolarQuote(db.Model):
    __tablename__ = 'solar_quote'
    id         = db.Column(db.Integer, primary_key=True)
    code       = db.Column(db.String(8), unique=True, nullable=False, index=True)
    data       = db.Column(db.Text, nullable=False)   # JSON
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


# ───────────────────────────── HELPERS ──────────────────────────────────── #

USD_TO_GHS_DEFAULT = 15.5

def get_ghs_rate():
    """Read the live GHS rate from DB; fall back to default."""
    try:
        row = SiteSettings.query.get('usd_to_ghs_rate')
        return float(row.value) if row else USD_TO_GHS_DEFAULT
    except Exception:
        return USD_TO_GHS_DEFAULT

# Keep a module-level alias for the email helper (uses live rate each call)
USD_TO_GHS = USD_TO_GHS_DEFAULT

def usd_to_ghs(usd):
    if usd is None:
        return None
    return usd * get_ghs_rate()

app.jinja_env.globals['usd_to_ghs'] = usd_to_ghs

def _ext(filename):
    return filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

def allowed(filename, kind='all'):
    e = _ext(filename)
    if kind == 'image': return e in ALLOWED_IMAGES
    if kind == 'video': return e in ALLOWED_VIDEOS
    return e in ALLOWED_ALL

def save_upload(file, subfolder):
    """Upload to Cloudinary when configured, otherwise save to local disk."""
    if USE_CLOUDINARY:
        resource_type = 'video' if _ext(file.filename) in ALLOWED_VIDEOS else 'image'
        result = cloudinary.uploader.upload(
            file,
            folder=f"dealership/{subfolder}",
            resource_type=resource_type,
        )
        # Store Cloudinary public_id prefixed with 'cld:' so we can delete it later
        return f"cld:{result['public_id']}|{result['secure_url']}"
    # Local fallback
    filename = secure_filename(file.filename)
    ext = _ext(filename)
    unique = f"{uuid.uuid4().hex}.{ext}"
    folder = os.path.join(app.config['UPLOAD_FOLDER'], subfolder)
    os.makedirs(folder, exist_ok=True)
    file.save(os.path.join(folder, unique))
    return f"uploads/{subfolder}/{unique}"


def media_url(path):
    """Return a URL for a stored media path (Cloudinary or local static)."""
    if path and path.startswith('cld:'):
        return path.split('|', 1)[1]   # Cloudinary secure_url
    return url_for('static', filename=path) if path else ''

app.jinja_env.globals['media_url'] = media_url


def _delete_media_file(path):
    """Delete a media file from Cloudinary or local disk."""
    if not path:
        return
    if path.startswith('cld:'):
        try:
            public_id = path.split('cld:', 1)[1].split('|')[0]
            cloudinary.uploader.destroy(public_id)
        except Exception:
            pass
    else:
        try:
            os.remove(os.path.join('static', path))
        except Exception:
            pass

def admin_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if 'admin_id' not in session:
            flash('Please log in to access the admin panel.', 'error')
            return redirect(url_for('admin_login'))
        return f(*args, **kwargs)
    return decorated

# ─── Mobile photo upload via QR code ────────────────────────────────────── #
MOBILE_SESSIONS = {}   # token → {'files': [...], 'subfolder': str}

@app.route('/admin/mobile-token', methods=['POST'])
@admin_required
def admin_mobile_token():
    data      = request.get_json(silent=True) or {}
    subfolder = data.get('subfolder', 'vehicles')
    if subfolder not in ('vehicles', 'solar'):
        subfolder = 'vehicles'
    token = uuid.uuid4().hex
    MOBILE_SESSIONS[token] = {'files': [], 'subfolder': subfolder}
    return jsonify({'token': token})

@app.route('/mobile-upload/<token>')
def mobile_upload_page(token):
    if token not in MOBILE_SESSIONS:
        return '<h2 style="font-family:sans-serif;text-align:center;padding:3rem;color:#dc2626;">Session expired — please scan the QR code again.</h2>', 403
    return render_template('admin/mobile_upload.html', token=token, business=BUSINESS)

@app.route('/mobile-upload/<token>', methods=['POST'])
def mobile_upload_receive(token):
    if token not in MOBILE_SESSIONS:
        return jsonify({'ok': False, 'error': 'Session expired'}), 403
    subfolder = MOBILE_SESSIONS[token]['subfolder']
    uploaded  = []
    for f in request.files.getlist('files'):
        if f and f.filename and allowed(f.filename, 'image'):
            path = save_upload(f, subfolder)
            if path:
                MOBILE_SESSIONS[token]['files'].append(path)
                uploaded.append({'path': path, 'url': media_url(path)})
    return jsonify({'ok': True, 'count': len(uploaded), 'uploaded': uploaded})

@app.route('/admin/mobile-poll/<token>')
@admin_required
def admin_mobile_poll(token):
    data  = MOBILE_SESSIONS.get(token, {'files': []})
    files = data.get('files', [])
    return jsonify({'files': files, 'urls': [media_url(f) for f in files]})

@app.template_filter('currency')
def currency_filter(value):
    if value is None: return 'N/A'
    return f"${value:,.2f}"

@app.template_filter('ghs')
def ghs_filter(value):
    if value is None: return 'N/A'
    return f"GH₵ {value * get_ghs_rate():,.2f}"

@app.template_filter('comma')
def comma_filter(value):
    if value is None: return 'N/A'
    return f"{int(value):,}"

@app.context_processor
def inject_globals():
    try:
        unread         = Contact.query.filter_by(status='unread').count()
        pending        = LoanApplication.query.filter_by(status='pending').count()
        pending_orders = CarOrder.query.filter_by(status='pending').count()
    except Exception:
        unread         = 0
        pending        = 0
        pending_orders = 0

    # Resolve referral SP for OG tags.
    # Check the live URL ?ref= first so social-media scrapers (which have no
    # session) still get the personalised preview card.
    ref_code = session.get('ref_sp_code', '')
    ref_name = session.get('ref_sp_name', '')
    url_ref  = request.args.get('ref', '').strip().upper()
    if url_ref and not ref_code:
        try:
            sp = Salesperson.query.filter_by(code=url_ref, active=True).first()
            if sp:
                ref_code = sp.code
                ref_name = sp.full_name
        except Exception:
            pass

    return {
        'business':        BUSINESS,
        'current_year':    datetime.utcnow().year,
        'unread_contacts': unread,
        'pending_loans':   pending,
        'pending_orders':  pending_orders,
        'USD_TO_GHS':      get_ghs_rate(),
        'ref_sp_code':     ref_code,
        'ref_sp_name':     ref_name,
    }


# ── Salesperson referral tracking ───────────────────────────────────────── #
@app.before_request
def capture_referral():
    """Store salesperson referral in session when ?ref=CODE is in the URL."""
    ref = request.args.get('ref', '').strip()
    if ref:
        sp = Salesperson.query.filter_by(code=ref.upper(), active=True).first() \
          or Salesperson.query.filter_by(code=ref, active=True).first()
        if sp:
            session['ref_sp_id']    = sp.id
            session['ref_sp_code']  = sp.code
            session['ref_sp_name']  = sp.full_name
            session['ref_sp_phone'] = sp.phone or ''
            flash(
                f'🎉 Welcome! {sp.full_name} has shared an exclusive deal with you. '
                f'Don\'t worry about financing — loan packages available via Republic Bank or Access Bank. '
                f'Mention code {sp.code} when you order to unlock your discount!',
                'success'
            )
            session['finance_flash_shown'] = True  # suppress the generic flash on home page

def _sp_session():
    """Return (sp_code, sp_name, sp_phone) from the current session."""
    return (
        session.get('ref_sp_code', ''),
        session.get('ref_sp_name', ''),
        session.get('ref_sp_phone', ''),
    )

def send_order_notification(order, vehicle):
    """Email order/test-drive details to admin."""
    if not SMTP_USER or not SMTP_PASSWORD:
        return
    try:
        sp_line = ''
        if order.sp_code:
            sp_line = f"\n  REFERRED BY\n  -----------\n  Salesperson : {order.sp_name}\n  Code        : {order.sp_code}\n  Phone       : {order.sp_phone}\n"
        otype = 'TEST DRIVE BOOKING' if order.order_type == 'test_drive' else 'VEHICLE ORDER'
        vprice = f"${vehicle.price:,.2f}"
        body  = f"""\
NEW {otype} — #{order.id}
{'=' * 54}

  VEHICLE
  -------
  {vehicle.name}{' (' + str(vehicle.year) + ')' if vehicle.year else ''}
  Price : {vprice}
{sp_line}
  CUSTOMER
  --------
  Name     : {order.name}
  Email    : {order.email}
  Phone    : {order.phone}
  Location : {order.location or 'Not provided'}
  {'Date : ' + order.preferred_date if order.preferred_date else ''}
  {'Time : ' + order.preferred_time if order.preferred_time else ''}
  Notes    : {order.notes or 'None'}

{'=' * 54}
Submitted : {order.created_at.strftime('%d %B %Y at %H:%M UTC')}
"""
        msg = MIMEMultipart('alternative')
        msg['Subject'] = f"[{otype} #{order.id}] {vehicle.name} — {order.name}"
        msg['From']    = SMTP_USER
        msg['To']      = ', '.join(NOTIFY_EMAILS)
        msg.attach(MIMEText(body, 'plain'))
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
            server.ehlo(); server.starttls()
            server.login(SMTP_USER, SMTP_PASSWORD)
            server.sendmail(SMTP_USER, NOTIFY_EMAILS, msg.as_string())
    except Exception as e:
        print(f'[EMAIL] Order notification failed: {e}')

def send_salesperson_welcome(sp):
    """Email the new salesperson their code + login details, and notify admin."""
    if not SMTP_USER or not SMTP_PASSWORD:
        return
    try:
        login_url   = 'https://chinacarsinghana.com/sales/login'
        ref_link    = f'https://chinacarsinghana.com?ref={sp.code}'
        joined_date = sp.created_at.strftime('%d %B %Y at %H:%M UTC')

        # ── Email to the salesperson ──────────────────────────────────────
        sp_body = f"""\
Welcome to the China Cars in Ghana Sales Team!
{'=' * 54}

Hi {sp.full_name},

Your account has been created. Here are your details:

  YOUR UNIQUE CODE : {sp.code}
  LOGIN URL        : {login_url}
  YOUR REFERRAL LINK:
  {ref_link}

HOW IT WORKS
------------
1. Share your referral link on WhatsApp, Facebook, or
   anywhere your audience is.
2. When a customer visits through your link and places
   an order or applies for a loan, it is automatically
   linked to you.
3. You earn 1% commission on the full sale value for
   every confirmed (fully paid) sale you bring in.

LOG IN TO YOUR DASHBOARD
-------------------------
Visit {login_url} and enter your code ({sp.code})
together with the password you set during registration.

Questions? Contact us at +233 50 356 6913 or reply to
this email.

Good luck and happy selling!
— China Cars in Ghana Team
{'=' * 54}
chinacarsinghana.com · Accra, Ghana
"""
        msg_sp = MIMEMultipart('alternative')
        msg_sp['Subject'] = f'Welcome to the Team — Your Code is {sp.code}'
        msg_sp['From']    = SMTP_USER
        msg_sp['To']      = sp.email
        msg_sp.attach(MIMEText(sp_body, 'plain'))

        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
            server.ehlo(); server.starttls()
            server.login(SMTP_USER, SMTP_PASSWORD)
            server.sendmail(SMTP_USER, [sp.email], msg_sp.as_string())

        # ── Admin notification ────────────────────────────────────────────
        admin_body = f"""\
NEW SALESPERSON REGISTRATION
{'=' * 54}

A new salesperson has self-registered via the public link.

  Name      : {sp.full_name}
  Code      : {sp.code}
  Email     : {sp.email or '—'}
  Phone     : {sp.phone or '—'}
  Commission: {sp.commission_pct}%
  Registered: {joined_date}

Their welcome email (with code + login details) has been
sent automatically to {sp.email or 'their email address'}.

Manage the team at:
https://chinacarsinghana.com/admin/salespersons
{'=' * 54}
"""
        msg_admin = MIMEMultipart('alternative')
        msg_admin['Subject'] = f'[New Salesperson] {sp.full_name} — Code {sp.code}'
        msg_admin['From']    = SMTP_USER
        msg_admin['To']      = ', '.join(NOTIFY_EMAILS)
        msg_admin.attach(MIMEText(admin_body, 'plain'))

        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
            server.ehlo(); server.starttls()
            server.login(SMTP_USER, SMTP_PASSWORD)
            server.sendmail(SMTP_USER, NOTIFY_EMAILS, msg_admin.as_string())

    except Exception as e:
        print(f'[EMAIL] Salesperson welcome failed: {e}')


def send_new_listing_alert(listing_type, name, price, listing_url, action='new'):
    """Email all active salespersons when a listing is added or updated."""
    if not SMTP_USER or not SMTP_PASSWORD:
        return
    try:
        with app.app_context():
            salespersons = Salesperson.query.filter_by(active=True).all()
            recipients   = [sp for sp in salespersons if sp.email]
            if not recipients:
                return

            emoji    = '☀️' if listing_type == 'solar' else '🚗'
            type_lbl = 'Solar System' if listing_type == 'solar' else 'Vehicle'
            is_new   = action == 'new'
            verb     = 'added to' if is_new else 'updated in'
            subj_pfx = f'{emoji} New {type_lbl}' if is_new else f'{emoji} Updated {type_lbl}'

            for sp in recipients:
                ref_link = f'https://chinacarsinghana.com?ref={sp.code}'
                body = f"""\
{emoji} {type_lbl.upper()} {'LISTING' if is_new else 'UPDATE'} — China Cars in Ghana
{'=' * 54}

Hi {sp.full_name},

A {type_lbl.lower()} has just been {verb} the inventory:

  {type_lbl.upper()} : {name}
  PRICE      : ${price:,.2f}
  VIEW LINK  : {listing_url}

Share it with your customers using your referral link
and earn 1% commission on every confirmed sale:

  {ref_link}

Log in to your dashboard:
  https://chinacarsinghana.com/sales/login

Good luck!
— China Cars in Ghana Team
{'=' * 54}
chinacarsinghana.com · Accra, Ghana
"""
                msg = MIMEMultipart('alternative')
                msg['Subject'] = f'{subj_pfx}: {name} — Start Sharing!'
                msg['From']    = SMTP_USER
                msg['To']      = sp.email
                msg.attach(MIMEText(body, 'plain'))

                try:
                    with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
                        server.ehlo(); server.starttls()
                        server.login(SMTP_USER, SMTP_PASSWORD)
                        server.sendmail(SMTP_USER, [sp.email], msg.as_string())
                except Exception as e:
                    print(f'[EMAIL] Listing alert to {sp.email} failed: {e}')

    except Exception as e:
        print(f'[EMAIL] send_new_listing_alert failed: {e}')


def send_price_change_alert(product_type, product_id, product_name, old_price_usd, new_price_usd):
    """Email all loan applicants for a product when its price changes."""
    if not SMTP_USER or not SMTP_PASSWORD:
        return
    # Only fire when price actually changed by more than a rounding artefact
    if abs(new_price_usd - old_price_usd) < 0.01:
        return
    try:
        with app.app_context():
            rate = get_ghs_rate()
            old_ghs = old_price_usd * rate
            new_ghs = new_price_usd * rate
            direction = 'decreased' if new_ghs < old_ghs else 'increased'

            apps = LoanApplication.query.filter_by(
                product_type=product_type,
                product_id=product_id,
            ).all()
            recipients = [a for a in apps if a.email]
            if not recipients:
                return

            loan_url = 'https://chinacarsinghana.com/loan'

            for a in recipients:
                body = f"""\
PRICE UPDATE — {product_name}
{'=' * 54}

Dear {a.first_name},

We are writing to let you know that the price of the
product you applied to finance has {direction}:

  PRODUCT   : {product_name}
  OLD PRICE : GH₵ {old_ghs:,.2f}
  NEW PRICE : GH₵ {new_ghs:,.2f}

Because the price has changed, your previous loan
application (Ref #{a.id}) may no longer reflect the
correct figures. We kindly ask you to submit a new
application with the updated price so that Republic Bank
can process it accurately.

Apply here:
  {loan_url}

If you have any questions, please contact us:
  kofi@chinacarsinghana.com

We appreciate your patience and look forward to helping
you complete your purchase.

— China Cars in Ghana Team
{'=' * 54}
chinacarsinghana.com · Accra, Ghana
"""
                msg = MIMEMultipart('alternative')
                msg['Subject'] = f'Price Update: {product_name} — Please Reapply for Financing'
                msg['From']    = SMTP_USER
                msg['To']      = a.email
                msg.attach(MIMEText(body, 'plain'))
                try:
                    with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
                        server.ehlo(); server.starttls()
                        server.login(SMTP_USER, SMTP_PASSWORD)
                        server.sendmail(SMTP_USER, [a.email], msg.as_string())
                except Exception as e:
                    print(f'[EMAIL] Price alert to {a.email} failed: {e}')
    except Exception as e:
        print(f'[EMAIL] send_price_change_alert failed: {e}')


def send_manual_price_alert(product_type, product_id, product_name, price_usd):
    """Manually notify all loan applicants of the current price for a product."""
    if not SMTP_USER or not SMTP_PASSWORD:
        return
    try:
        with app.app_context():
            rate    = get_ghs_rate()
            ghs     = price_usd * rate
            apps    = LoanApplication.query.filter_by(
                product_type=product_type, product_id=product_id
            ).all()
            recipients = [a for a in apps if a.email]
            if not recipients:
                return
            loan_url = 'https://chinacarsinghana.com/loan'
            for a in recipients:
                body = f"""\
IMPORTANT: UPDATED PRICE — {product_name}
{'=' * 54}

Dear {a.first_name},

We are reaching out regarding your loan application
(Ref #{a.id}) for the following product:

  PRODUCT       : {product_name}
  CURRENT PRICE : GH₵ {ghs:,.2f}

To ensure your financing is processed at the correct
amount, please submit a fresh application using the
link below:

  {loan_url}

If you have any questions, contact us at:
  kofi@chinacarsinghana.com

Thank you for choosing China Cars in Ghana.

— China Cars in Ghana Team
{'=' * 54}
chinacarsinghana.com · Accra, Ghana
"""
                msg = MIMEMultipart('alternative')
                msg['Subject'] = f'Action Required: Please Reapply — {product_name}'
                msg['From']    = SMTP_USER
                msg['To']      = a.email
                msg.attach(MIMEText(body, 'plain'))
                try:
                    with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
                        server.ehlo(); server.starttls()
                        server.login(SMTP_USER, SMTP_PASSWORD)
                        server.sendmail(SMTP_USER, [a.email], msg.as_string())
                except Exception as e:
                    print(f'[EMAIL] Manual price alert to {a.email} failed: {e}')
    except Exception as e:
        print(f'[EMAIL] send_manual_price_alert failed: {e}')


# ────────────────────────── PUBLIC ROUTES ───────────────────────────────── #

@app.route('/static/img/og-sales.jpg')
def og_sales_image():
    """OG share image for the salesperson recruitment / registration page."""
    try:
        from PIL import Image, ImageDraw, ImageFont

        W, H = 1200, 630
        img  = Image.new('RGB', (W, H), color=(15, 23, 42))
        draw = ImageDraw.Draw(img)

        # Gold gradient bar at top
        for x in range(W):
            t = x / W
            r = int(234 + (16  - 234) * t)
            g = int(179 + (185 - 179) * t)
            b = int(8   + (129 - 8)   * t)
            draw.line([(x, 0), (x, 10)], fill=(r, g, b))

        # Subtle corner glow (top-right, blue)
        for r in range(280, 0, -20):
            alpha = max(0, 18 - (280 - r) // 15)
            draw.ellipse([W - r, -r // 2, W + r // 2, r], fill=(37, 99, 235))

        try:
            f_hero  = ImageFont.load_default(size=88)
            f_big   = ImageFont.load_default(size=54)
            f_mid   = ImageFont.load_default(size=32)
            f_small = ImageFont.load_default(size=26)
            f_tag   = ImageFont.load_default(size=28)
        except Exception:
            f_hero = f_big = f_mid = f_small = f_tag = ImageFont.load_default()

        # Top label
        draw.text((80, 30), 'CHINA CARS IN GHANA', font=f_small, fill=(148, 163, 184))

        # Hero headline
        draw.text((80, 80),  'JOIN OUR',    font=f_hero, fill=(255, 255, 255))
        draw.text((80, 175), 'SALES TEAM',  font=f_hero, fill=(234, 179, 8))

        # Divider
        draw.rectangle([80, 280, 320, 286], fill=(234, 179, 8))

        # Commission pill (green)
        def pill(x, y, text, bg, fg, radius=26):
            tw = len(text) * 16 + 44
            draw.rounded_rectangle([x, y, x + tw, y + 54], radius=radius, fill=bg)
            draw.text((x + 22, y + 13), text, font=f_tag, fill=fg)
            return x + tw + 16

        nx = pill(80, 304, '1% COMMISSION',      (20, 83, 45),   (134, 239, 172))
        nx = pill(nx, 304, 'VEHICLES & SOLAR',  (30, 58, 138),  (147, 197, 253))
        pill(nx,  304, 'FREE TO JOIN',           (88, 28, 135),  (216, 180, 254))

        # White info card
        draw.rounded_rectangle([80, 378, W - 80, 530], radius=16,
                                fill=(255, 255, 255), outline=(234, 179, 8), width=2)
        draw.text((120, 396), 'Earn on every vehicle & solar system you sell',
                  font=f_mid, fill=(15, 23, 42))
        draw.text((120, 438), 'Get your personal referral link · Track earnings in real time',
                  font=f_small, fill=(71, 85, 105))
        draw.text((120, 476), 'Commission paid on full customer payment · Sign up free today',
                  font=f_small, fill=(71, 85, 105))

        # Footer
        draw.text((80,  554), 'chinacarsinghana.com/sales/register', font=f_mid,   fill=(96, 165, 250))
        draw.text((80,  598), 'Accra, Ghana',                         font=f_small, fill=(71, 85, 105))

        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=93)
        buf.seek(0)
        return send_file(buf, mimetype='image/jpeg',
                         max_age=86400,
                         as_attachment=False,
                         download_name='og-sales.jpg')
    except Exception as e:
        print(f'[OG SALES IMAGE] Error: {e}')
        from flask import abort
        abort(404)


@app.route('/static/img/og-default.jpg')
def og_default_image():
    """Dynamically generate the default OG share image (1200×630 JPEG)."""
    try:
        from PIL import Image, ImageDraw, ImageFont

        W, H = 1200, 630
        img  = Image.new('RGB', (W, H), color=(15, 23, 42))
        draw = ImageDraw.Draw(img)

        # Blue-to-orange gradient bar at top (8px)
        for x in range(W):
            t = x / W
            r = int(37  + (249 - 37)  * t)
            g = int(99  + (115 - 99)  * t)
            b = int(235 + (22  - 235) * t)
            draw.line([(x, 0), (x, 8)], fill=(r, g, b))

        # Subtle diagonal highlight
        for i in range(0, 300, 40):
            draw.line([(W - i, 0), (W, i)], fill=(255, 255, 255, 8), width=1)

        # Load scalable default fonts (Pillow 10+)
        try:
            f_big   = ImageFont.load_default(size=72)
            f_mid   = ImageFont.load_default(size=34)
            f_small = ImageFont.load_default(size=26)
            f_tag   = ImageFont.load_default(size=28)
        except Exception:
            f_big = f_mid = f_small = f_tag = ImageFont.load_default()

        # Business name
        draw.text((80, 80),  'CHINA CARS IN GHANA', font=f_big,   fill=(255, 255, 255))
        # Tagline
        draw.text((84, 175), 'Drive Excellence.  Power the Future.',  font=f_mid,   fill=(148, 163, 184))

        # Divider line
        draw.rectangle([80, 230, 400, 235], fill=(37, 99, 235))

        # Tag pills
        def pill(x, y, text, bg, fg):
            tw = len(text) * 17 + 40
            draw.rounded_rectangle([x, y, x + tw, y + 52], radius=26, fill=bg)
            draw.text((x + 20, y + 12), text, font=f_tag, fill=fg)

        pill(80,  260, 'VEHICLES',     (30, 58, 138),  (147, 197, 253))
        pill(300, 260, 'SOLAR',        (124, 45, 18),  (253, 186, 116))
        pill(470, 260, 'EASY FINANCE', (20, 83, 45),   (134, 239, 172))

        # Big price / value prop block
        draw.rectangle([80, 340, W - 80, 490], fill=(255, 255, 255, 0))
        draw.rounded_rectangle([80, 340, W - 80, 500], radius=16,
                                fill=(255, 255, 255), outline=(37, 99, 235), width=2)
        draw.text((120, 360), 'Premium imported vehicles & solar energy systems',
                  font=f_mid, fill=(15, 23, 42))
        draw.text((120, 405), 'Flexible financing via Republic Bank · Delivery across Ghana',
                  font=f_small, fill=(71, 85, 105))
        draw.text((120, 445), 'Call / WhatsApp:  +233 50 356 6913',
                  font=f_mid, fill=(37, 99, 235))

        # Domain footer
        draw.text((80,  540), 'chinacarsinghana.com',  font=f_mid,   fill=(96, 165, 250))
        draw.text((80,  585), 'Accra, Ghana',          font=f_small, fill=(71, 85, 105))

        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=93)
        buf.seek(0)
        return send_file(buf, mimetype='image/jpeg',
                         max_age=86400,
                         as_attachment=False,
                         download_name='og-default.jpg')
    except Exception as e:
        print(f'[OG IMAGE] Error: {e}')
        from flask import abort
        abort(404)

@app.route('/')
def index():
    if not session.get('finance_flash_shown'):
        flash(
            '💳 Don\'t let financing hold you back — we offer flexible loan packages '
            'through Republic Bank. Drive away today, pay comfortably over time. '
            'Apply in minutes!',
            'info'
        )
        session['finance_flash_shown'] = True

    featured_vehicles = Vehicle.query.filter_by(featured=True, available=True).limit(4).all()
    featured_solar    = SolarSystem.query.filter_by(featured=True, available=True).limit(2).all()
    latest_vehicles   = Vehicle.query.filter_by(available=True).order_by(Vehicle.created_at.desc()).limit(6).all()
    latest_solar      = SolarSystem.query.filter_by(available=True).order_by(SolarSystem.created_at.desc()).limit(4).all()
    newest_vehicle    = Vehicle.query.filter_by(available=True).order_by(Vehicle.created_at.desc()).first()
    return render_template('index.html',
                           featured_vehicles=featured_vehicles,
                           featured_solar=featured_solar,
                           latest_vehicles=latest_vehicles,
                           latest_solar=latest_solar,
                           newest_vehicle=newest_vehicle)


@app.route('/vehicles')
def vehicles():
    category  = request.args.get('category', '')
    fuel_type = request.args.get('fuel', '')
    min_price = request.args.get('min', type=float)
    max_price = request.args.get('max', type=float)
    search    = request.args.get('q', '').strip()
    sort      = request.args.get('sort', 'newest')

    q = Vehicle.query.filter_by(available=True)
    if category:  q = q.filter_by(category=category)
    if fuel_type: q = q.filter_by(fuel_type=fuel_type)
    if min_price: q = q.filter(Vehicle.price >= min_price)
    if max_price: q = q.filter(Vehicle.price <= max_price)
    if search:    q = q.filter(Vehicle.name.ilike(f'%{search}%'))

    if sort == 'price_asc':  q = q.order_by(Vehicle.price.asc())
    elif sort == 'price_desc': q = q.order_by(Vehicle.price.desc())
    else: q = q.order_by(Vehicle.created_at.desc())

    vehicle_list = q.all()
    categories = [r[0] for r in db.session.query(Vehicle.category).distinct() if r[0]]
    fuels      = [r[0] for r in db.session.query(Vehicle.fuel_type).distinct() if r[0]]

    return render_template('vehicles.html',
                           vehicles=vehicle_list,
                           categories=categories,
                           fuels=fuels,
                           current_category=category,
                           current_fuel=fuel_type,
                           current_sort=sort,
                           search=search)


@app.route('/vehicles/<int:vid>')
def vehicle_detail(vid):
    vehicle = Vehicle.query.get_or_404(vid)
    related = Vehicle.query.filter(Vehicle.id != vid,
                                   Vehicle.available == True,
                                   Vehicle.category == vehicle.category).limit(4).all()
    return render_template('vehicle_detail.html', vehicle=vehicle, related=related)


@app.route('/solar')
def solar():
    category = request.args.get('category', '')
    q = SolarSystem.query.filter_by(available=True)
    if category: q = q.filter_by(category=category)
    solar_list = q.order_by(SolarSystem.created_at.desc()).all()
    categories = [r[0] for r in db.session.query(SolarSystem.category).distinct() if r[0]]
    return render_template('solar.html',
                           solar_systems=solar_list,
                           categories=categories,
                           current_category=category)


@app.route('/solar/<int:sid>')
def solar_detail(sid):
    solar   = SolarSystem.query.get_or_404(sid)
    related = SolarSystem.query.filter(SolarSystem.id != sid,
                                        SolarSystem.available == True).limit(4).all()
    return render_template('solar_detail.html', solar=solar, related=related)


@app.route('/loan', methods=['GET', 'POST'])
def loan():
    all_vehicles = Vehicle.query.filter_by(available=True).all()
    all_solar    = SolarSystem.query.filter_by(available=True).all()

    if request.method == 'POST':
        product_type  = request.form.get('product_type')
        product_name  = ''
        product_price = 0.0
        product_id    = None

        if product_type == 'vehicle':
            product_id = request.form.get('vehicle_id', type=int)
            if product_id:
                v = Vehicle.query.get(product_id)
                if v:
                    product_name  = f"{v.name}{' (' + str(v.year) + ')' if v.year else ''}"
                    product_price = v.price
        elif product_type == 'solar':
            product_id = request.form.get('solar_id', type=int)
            if product_id:
                s = SolarSystem.query.get(product_id)
                if s:
                    product_name  = s.name
                    product_price = s.price

        sp_code, sp_name, sp_phone = _sp_session()
        app_obj = LoanApplication(
            first_name        = request.form.get('first_name', '').strip(),
            last_name         = request.form.get('last_name', '').strip(),
            email             = request.form.get('email', '').strip(),
            phone             = request.form.get('phone', '').strip(),
            date_of_birth     = request.form.get('date_of_birth'),
            id_number         = request.form.get('id_number', '').strip(),
            address           = request.form.get('address', '').strip(),
            product_type      = product_type,
            product_id        = product_id,
            product_name      = product_name,
            product_price     = product_price,
            loan_amount       = request.form.get('loan_amount', type=float),
            deposit_amount    = request.form.get('deposit_amount', type=float),
            loan_term         = request.form.get('loan_term', type=int),
            employment_status = request.form.get('employment_status'),
            employer          = request.form.get('employer', '').strip(),
            annual_income     = request.form.get('annual_income', type=float),
            message           = request.form.get('message', '').strip(),
            bank_choice       = request.form.get('bank_choice', 'republic'),
            sp_code           = sp_code,
            sp_name           = sp_name,
            sp_phone          = sp_phone,
        )
        db.session.add(app_obj)
        db.session.commit()
        send_loan_notification(app_obj)
        flash('Your loan application has been submitted successfully! We will contact you within 24–48 hours.', 'success')
        return redirect(url_for('loan'))

    preselect_v    = request.args.get('vehicle', type=int)
    preselect_s    = request.args.get('solar', type=int)
    preselect_type = request.args.get('product_type', '')  # 'vehicle' or 'solar'
    if preselect_type == 'solar' and not preselect_s:
        preselect_s = -1  # flag to pre-check solar radio without a specific item
    loan_rate_row = SiteSettings.query.filter_by(key='loan_interest_rate').first()
    loan_interest_rate = float(loan_rate_row.value) if loan_rate_row else 7.0
    return render_template('loan.html',
                           vehicles=all_vehicles,
                           solar_systems=all_solar,
                           preselect_vehicle=preselect_v,
                           preselect_solar=preselect_s,
                           loan_interest_rate=loan_interest_rate)


@app.route('/contact', methods=['GET', 'POST'])
def contact():
    if request.method == 'POST':
        sp_code, sp_name, sp_phone = _sp_session()
        c = Contact(
            name     = request.form.get('name', '').strip(),
            email    = request.form.get('email', '').strip(),
            phone    = request.form.get('phone', '').strip(),
            subject  = request.form.get('subject', '').strip(),
            message  = request.form.get('message', '').strip(),
            sp_code  = sp_code,
            sp_name  = sp_name,
            sp_phone = sp_phone,
        )
        db.session.add(c)
        db.session.commit()
        flash('Thank you! Your message has been sent. We will be in touch shortly.', 'success')
        return redirect(url_for('contact'))
    return render_template('contact.html')


@app.route('/about')
def about():
    return render_template('about.html')


@app.route('/order/<int:vid>', methods=['GET', 'POST'])
def order_vehicle(vid):
    vehicle = Vehicle.query.get_or_404(vid)
    if request.method == 'POST':
        sp_code, sp_name, sp_phone = _sp_session()
        order = CarOrder(
            order_type     = request.form.get('order_type', 'order'),
            vehicle_id     = vid,
            name           = request.form.get('name', '').strip(),
            email          = request.form.get('email', '').strip(),
            phone          = request.form.get('phone', '').strip(),
            location       = request.form.get('location', '').strip(),
            latitude       = request.form.get('latitude', type=float),
            longitude      = request.form.get('longitude', type=float),
            preferred_date = request.form.get('preferred_date'),
            preferred_time = request.form.get('preferred_time'),
            notes          = request.form.get('notes', '').strip(),
            sp_code        = sp_code,
            sp_name        = sp_name,
            sp_phone       = sp_phone,
        )
        db.session.add(order)
        db.session.commit()
        send_order_notification(order, vehicle)
        order_type_label = 'test drive' if order.order_type == 'test_drive' else 'order'
        flash(
            f'Your {order_type_label} request for the {vehicle.name} has been submitted! '
            f'We will contact you within 24 hours. If you do not hear from us, please email us directly at kofi@chinacarsinghana.com.',
            'success'
        )
        return redirect(url_for('vehicle_detail', vid=vid))
    order_type = request.args.get('type', 'order')  # 'order' or 'test_drive'
    return render_template('order.html', vehicle=vehicle, order_type=order_type)


@app.route('/admin/orders')
@admin_required
def admin_orders():
    order_type = request.args.get('type', '')
    status     = request.args.get('status', '')
    q = CarOrder.query
    if order_type: q = q.filter_by(order_type=order_type)
    if status:     q = q.filter_by(status=status)
    orders = q.order_by(CarOrder.created_at.desc()).all()
    return render_template('admin/orders.html', orders=orders, current_type=order_type, current_status=status)


@app.route('/admin/orders/<int:oid>/status', methods=['POST'])
@admin_required
def admin_update_order_status(oid):
    order = CarOrder.query.get_or_404(oid)
    order.status = request.form.get('status', 'pending')
    db.session.commit()
    flash('Order status updated.', 'success')
    return redirect(url_for('admin_orders'))


@app.route('/admin/orders/test-drive-blast', methods=['POST'])
@admin_required
def admin_test_drive_blast():
    """Email every test drive customer telling them to contact their sales agent for showroom directions."""
    if not SMTP_USER or not SMTP_PASSWORD:
        flash('SMTP not configured — no emails sent.', 'error')
        return redirect(url_for('admin_orders'))

    orders = CarOrder.query.filter_by(order_type='test_drive').all()
    if not orders:
        flash('No test drive bookings found.', 'error')
        return redirect(url_for('admin_orders'))

    SHOWROOM_LINK = 'https://maps.app.goo.gl/M2xCkDrKbhEUrqyn8?g_st=iwb'
    DEALERSHIP_PHONE = '+233 50 356 6913'
    sent = 0
    failed = 0

    def _send_all(app_ctx, order_list):
        nonlocal sent, failed
        with app_ctx:
            for o in order_list:
                if not o.email:
                    continue
                vehicle_name = o.vehicle.name if o.vehicle else 'the vehicle'
                if o.sp_name and o.sp_phone:
                    contact_block = (
                        f"Please reach out directly to your sales agent:\n\n"
                        f"  Agent : {o.sp_name}\n"
                        f"  Phone : {o.sp_phone}\n\n"
                        f"They will guide you on the best time to visit and answer any questions."
                    )
                elif o.sp_name:
                    contact_block = (
                        f"Please reach out directly to your sales agent, {o.sp_name}, "
                        f"or call us on {DEALERSHIP_PHONE} to arrange your visit."
                    )
                else:
                    contact_block = (
                        f"Please call us on {DEALERSHIP_PHONE} or reply to this email "
                        f"to arrange a convenient time for your visit."
                    )

                body = f"""\
Hi {o.name},

Thank you for booking a test drive for the {vehicle_name}!

We are excited to have you visit our showroom in Tema so you can experience the vehicle in person.

──────────────────────────────────────
SHOWROOM LOCATION — TEMA
──────────────────────────────────────
Open in Google Maps:
{SHOWROOM_LINK}

──────────────────────────────────────
NEXT STEP
──────────────────────────────────────
{contact_block}

We look forward to seeing you!

──────────────────────────────────────
China Cars in Ghana
{DEALERSHIP_PHONE}
kofi@chinacarsinghana.com
https://chinacarsinghana.com
"""
                try:
                    msg = MIMEMultipart('alternative')
                    msg['Subject'] = f'Your Test Drive — Visit Us at Our Tema Showroom'
                    msg['From']    = SMTP_USER
                    msg['To']      = o.email
                    msg.attach(MIMEText(body, 'plain'))
                    with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
                        server.ehlo(); server.starttls()
                        server.login(SMTP_USER, SMTP_PASSWORD)
                        server.sendmail(SMTP_USER, [o.email], msg.as_string())
                    sent += 1
                except Exception as e:
                    print(f'[EMAIL] Test drive blast failed for {o.email}: {e}')
                    failed += 1

    threading.Thread(target=_send_all, args=(app.app_context(), orders), daemon=True).start()

    total = len([o for o in orders if o.email])
    flash(
        f'Sending showroom location email to {total} test drive customer{"s" if total != 1 else ""}. '
        f'Delivery may take a minute.',
        'success'
    )
    return redirect(url_for('admin_orders'))


# ────────────────────────── ADMIN ROUTES ────────────────────────────────── #

@app.route('/admin/login', methods=['GET', 'POST'])
def admin_login():
    if 'admin_id' in session:
        return redirect(url_for('admin_dashboard'))
    if request.method == 'POST':
        user = AdminUser.query.filter_by(username=request.form.get('username')).first()
        if user and user.check_password(request.form.get('password', '')):
            session['admin_id']   = user.id
            session['admin_name'] = user.full_name or user.username
            return redirect(url_for('admin_dashboard'))
        flash('Invalid username or password.', 'error')
    return render_template('admin/login.html')


@app.route('/admin/logout')
def admin_logout():
    session.clear()
    return redirect(url_for('admin_login'))


@app.route('/admin')
@app.route('/admin/dashboard')
@admin_required
def admin_dashboard():
    # Aggregate commission/revenue in two SQL queries (avoids N+1 per salesperson)
    from sqlalchemy import func as sqlfunc
    total_sp_revenue = db.session.query(
        sqlfunc.sum(SalespersonSale.sale_amount)
    ).filter(SalespersonSale.status != 'cancelled').scalar() or 0

    total_sp_commission = db.session.query(
        sqlfunc.sum(SalespersonSale.sale_amount * Salesperson.commission_pct / 100.0)
    ).join(Salesperson, SalespersonSale.salesperson_id == Salesperson.id
    ).filter(SalespersonSale.status != 'cancelled').scalar() or 0

    rate = get_ghs_rate()

    stats = {
        'vehicles':          Vehicle.query.count(),
        'solar':             SolarSystem.query.count(),
        'applications':      LoanApplication.query.count(),
        'pending':           LoanApplication.query.filter_by(status='pending').count(),
        'contacts':          Contact.query.filter_by(status='unread').count(),
        'approved':          LoanApplication.query.filter_by(status='approved').count(),
        'salespersons':      Salesperson.query.filter_by(active=True).count(),
        'sp_loans':          LoanApplication.query.filter(LoanApplication.sp_code != None, LoanApplication.sp_code != '').count(),
        'sp_orders':         CarOrder.query.filter(CarOrder.sp_code != None, CarOrder.sp_code != '').count(),
        'total_sp_revenue':  total_sp_revenue,
        'total_commission':  total_sp_commission,
        'total_commission_ghs': total_sp_commission * rate,
        'total_revenue_ghs': total_sp_revenue * rate,
    }
    recent_apps      = LoanApplication.query.order_by(LoanApplication.created_at.desc()).limit(6).all()
    recent_contacts  = Contact.query.order_by(Contact.created_at.desc()).limit(5).all()

    # Salesperson leaderboard — highest revenue first, then A-Z
    salespersons = sorted(
        Salesperson.query.filter_by(active=True).all(),
        key=lambda sp: (-sp.total_revenue(), sp.full_name.lower())
    )

    # Recent salesperson-attributed activity (loans + orders merged, last 10)
    sp_loan_activity = (LoanApplication.query
        .filter(LoanApplication.sp_code != None, LoanApplication.sp_code != '')
        .order_by(LoanApplication.created_at.desc()).limit(20).all())
    sp_order_activity = (CarOrder.query
        .filter(CarOrder.sp_code != None, CarOrder.sp_code != '')
        .order_by(CarOrder.created_at.desc()).limit(20).all())

    sp_activity = []
    for a in sp_loan_activity:
        sp_activity.append({
            'type': 'loan', 'sp_code': a.sp_code, 'sp_name': a.sp_name,
            'customer': f"{a.first_name} {a.last_name}",
            'product': a.product_name or a.product_type,
            'amount': a.loan_amount, 'status': a.status,
            'date': a.created_at,
        })
    for o in sp_order_activity:
        sp_activity.append({
            'type': 'order', 'sp_code': o.sp_code, 'sp_name': o.sp_name,
            'customer': o.name,
            'product': (o.vehicle.name if o.vehicle else 'Vehicle') + (' (Test Drive)' if o.order_type == 'test_drive' else ''),
            'amount': o.vehicle.price if o.vehicle else None, 'status': o.status,
            'date': o.created_at,
        })
    sp_activity.sort(key=lambda x: x['date'], reverse=True)
    sp_activity = sp_activity[:15]

    return render_template('admin/dashboard.html',
                           stats=stats,
                           recent_apps=recent_apps,
                           recent_contacts=recent_contacts,
                           salespersons=salespersons,
                           sp_activity=sp_activity)


# ── Vehicles ── #

@app.route('/admin/vehicles')
@admin_required
def admin_vehicles():
    q = request.args.get('q', '').strip()
    query = Vehicle.query
    if q:
        like = f'%{q}%'
        query = query.filter(
            db.or_(
                Vehicle.name.ilike(like),
                Vehicle.category.ilike(like),
                Vehicle.fuel_type.ilike(like),
                Vehicle.tagline.ilike(like),
            )
        )
    vlist = query.order_by(Vehicle.created_at.desc()).all()
    return render_template('admin/vehicles.html', vehicles=vlist, search_q=q)


@app.route('/admin/vehicles/add', methods=['GET', 'POST'])
@admin_required
def admin_add_vehicle():
    if request.method == 'POST':
        feats = [f.strip() for f in request.form.get('features', '').splitlines() if f.strip()]
        _price_ghs = request.form.get('price_ghs', type=float) or 0
        v = Vehicle(
            name         = request.form.get('name'),
            tagline      = request.form.get('tagline'),
            description  = request.form.get('description'),
            price        = _price_ghs / get_ghs_rate() if _price_ghs else 0,
            category     = request.form.get('category'),
            year         = request.form.get('year', type=int),
            mileage      = request.form.get('mileage', type=int),
            transmission = request.form.get('transmission'),
            fuel_type    = request.form.get('fuel_type'),
            engine       = request.form.get('engine'),
            color        = request.form.get('color'),
            seats        = request.form.get('seats', type=int),
            features     = json.dumps(feats),
            images       = '[]',
            videos       = '[]',
            featured     = 'featured' in request.form,
            available    = 'available' in request.form,
        )
        db.session.add(v)
        db.session.flush()

        imgs = []
        for f in request.files.getlist('images'):
            if f and f.filename and allowed(f.filename, 'image'):
                imgs.append(save_upload(f, 'vehicles'))
        mt = request.form.get('mobile_token', '')
        if mt and mt in MOBILE_SESSIONS:
            imgs.extend(MOBILE_SESSIONS.pop(mt)['files'])
        v.images = json.dumps(imgs)

        vids = []
        for f in request.files.getlist('videos'):
            if f and f.filename and allowed(f.filename, 'video'):
                vids.append(save_upload(f, 'vehicles'))
        v.videos = json.dumps(vids)

        db.session.commit()
        listing_url = f'https://chinacarsinghana.com/vehicles/{v.id}'
        threading.Thread(
            target=send_new_listing_alert,
            args=('vehicle', v.name, v.price, listing_url),
            daemon=True
        ).start()
        flash(f'Vehicle "{v.name}" added successfully!', 'success')
        return redirect(url_for('admin_vehicles'))
    return render_template('admin/vehicle_form.html', vehicle=None, action='Add')


@app.route('/admin/vehicles/edit/<int:vid>', methods=['GET', 'POST'])
@admin_required
def admin_edit_vehicle(vid):
    v = Vehicle.query.get_or_404(vid)
    if request.method == 'POST':
        feats = [f.strip() for f in request.form.get('features', '').splitlines() if f.strip()]
        _price_ghs = request.form.get('price_ghs', type=float) or 0
        _old_price = v.price
        v.name         = request.form.get('name')
        v.tagline      = request.form.get('tagline')
        v.description  = request.form.get('description')
        v.price        = _price_ghs / get_ghs_rate() if _price_ghs else 0
        v.category     = request.form.get('category')
        v.year         = request.form.get('year', type=int)
        v.mileage      = request.form.get('mileage', type=int)
        v.transmission = request.form.get('transmission')
        v.fuel_type    = request.form.get('fuel_type')
        v.engine       = request.form.get('engine')
        v.color        = request.form.get('color')
        v.seats        = request.form.get('seats', type=int)
        v.features     = json.dumps(feats)
        v.featured     = 'featured' in request.form
        v.available    = 'available' in request.form
        v.updated_at   = datetime.utcnow()

        existing_imgs = v.get_images()
        for f in request.files.getlist('images'):
            if f and f.filename and allowed(f.filename, 'image'):
                existing_imgs.append(save_upload(f, 'vehicles'))
        mt = request.form.get('mobile_token', '')
        if mt and mt in MOBILE_SESSIONS:
            existing_imgs.extend(MOBILE_SESSIONS.pop(mt)['files'])
        v.images = json.dumps(existing_imgs)

        existing_vids = v.get_videos()
        for f in request.files.getlist('videos'):
            if f and f.filename and allowed(f.filename, 'video'):
                existing_vids.append(save_upload(f, 'vehicles'))
        v.videos = json.dumps(existing_vids)

        # Keep salesperson sale records in sync with the updated price
        SalespersonSale.query.filter_by(product_type='vehicle', product_id=v.id).update(
            {'sale_amount': v.price}, synchronize_session=False
        )

        db.session.commit()
        listing_url = f'https://chinacarsinghana.com/vehicles/{v.id}'
        threading.Thread(
            target=send_new_listing_alert,
            args=('vehicle', v.name, v.price, listing_url, 'updated'),
            daemon=True
        ).start()
        threading.Thread(
            target=send_price_change_alert,
            args=('vehicle', v.id, v.name, _old_price, v.price),
            daemon=True
        ).start()
        flash(f'Vehicle "{v.name}" updated successfully!', 'success')
        return redirect(url_for('admin_vehicles'))
    return render_template('admin/vehicle_form.html', vehicle=v, action='Edit')


@app.route('/admin/vehicles/<int:vid>/notify-applicants', methods=['POST'])
@admin_required
def admin_notify_vehicle_applicants(vid):
    v = Vehicle.query.get_or_404(vid)
    threading.Thread(
        target=send_manual_price_alert,
        args=('vehicle', v.id, v.name, v.price),
        daemon=True
    ).start()
    flash(f'Price alert sent to all loan applicants for "{v.name}".', 'success')
    return redirect(url_for('admin_edit_vehicle', vid=vid))


@app.route('/admin/vehicles/delete/<int:vid>', methods=['POST'])
@admin_required
def admin_delete_vehicle(vid):
    v = Vehicle.query.get_or_404(vid)
    db.session.delete(v)
    db.session.commit()
    flash('Vehicle deleted.', 'success')
    return redirect(url_for('admin_vehicles'))


@app.route('/admin/vehicles/<int:vid>/delete-media', methods=['POST'])
@admin_required
def admin_delete_vehicle_media(vid):
    v    = Vehicle.query.get_or_404(vid)
    path = request.json.get('path')
    kind = request.json.get('kind', 'image')
    if kind == 'image':
        items = [i for i in v.get_images() if i != path]
        v.images = json.dumps(items)
    else:
        items = [i for i in v.get_videos() if i != path]
        v.videos = json.dumps(items)
    db.session.commit()
    _delete_media_file(path)
    return jsonify({'ok': True})


# ── Solar ── #

@app.route('/admin/solar')
@admin_required
def admin_solar():
    slist = SolarSystem.query.order_by(SolarSystem.created_at.desc()).all()
    return render_template('admin/solar.html', solar_systems=slist)


@app.route('/admin/solar/add', methods=['GET', 'POST'])
@admin_required
def admin_add_solar():
    if request.method == 'POST':
        feats = [f.strip() for f in request.form.get('features', '').splitlines() if f.strip()]
        _price_ghs = request.form.get('price_ghs', type=float) or 0
        s = SolarSystem(
            name             = request.form.get('name'),
            tagline          = request.form.get('tagline'),
            description      = request.form.get('description'),
            price            = _price_ghs / get_ghs_rate() if _price_ghs else 0,
            category         = request.form.get('category'),
            power_output     = request.form.get('power_output'),
            panel_count      = request.form.get('panel_count', type=int),
            battery_capacity = request.form.get('battery_capacity'),
            warranty         = request.form.get('warranty'),
            efficiency       = request.form.get('efficiency'),
            features         = json.dumps(feats),
            images           = '[]',
            videos           = '[]',
            featured         = 'featured' in request.form,
            available        = 'available' in request.form,
        )
        db.session.add(s)
        db.session.flush()

        imgs = []
        for f in request.files.getlist('images'):
            if f and f.filename and allowed(f.filename, 'image'):
                imgs.append(save_upload(f, 'solar'))
        mt = request.form.get('mobile_token', '')
        if mt and mt in MOBILE_SESSIONS:
            imgs.extend(MOBILE_SESSIONS.pop(mt)['files'])
        s.images = json.dumps(imgs)

        vids = []
        for f in request.files.getlist('videos'):
            if f and f.filename and allowed(f.filename, 'video'):
                vids.append(save_upload(f, 'solar'))
        s.videos = json.dumps(vids)

        db.session.commit()
        listing_url = f'https://chinacarsinghana.com/solar/{s.id}'
        threading.Thread(
            target=send_new_listing_alert,
            args=('solar', s.name, s.price, listing_url),
            daemon=True
        ).start()
        flash(f'Solar system "{s.name}" added successfully!', 'success')
        return redirect(url_for('admin_solar'))
    return render_template('admin/solar_form.html', solar=None, action='Add')


@app.route('/admin/solar/edit/<int:sid>', methods=['GET', 'POST'])
@admin_required
def admin_edit_solar(sid):
    s = SolarSystem.query.get_or_404(sid)
    if request.method == 'POST':
        feats = [f.strip() for f in request.form.get('features', '').splitlines() if f.strip()]
        _price_ghs = request.form.get('price_ghs', type=float) or 0
        _old_price = s.price
        s.name             = request.form.get('name')
        s.tagline          = request.form.get('tagline')
        s.description      = request.form.get('description')
        s.price            = _price_ghs / get_ghs_rate() if _price_ghs else 0
        s.category         = request.form.get('category')
        s.power_output     = request.form.get('power_output')
        s.panel_count      = request.form.get('panel_count', type=int)
        s.battery_capacity = request.form.get('battery_capacity')
        s.warranty         = request.form.get('warranty')
        s.efficiency       = request.form.get('efficiency')
        s.features         = json.dumps(feats)
        s.featured         = 'featured' in request.form
        s.available        = 'available' in request.form
        s.updated_at       = datetime.utcnow()

        existing_imgs = s.get_images()
        for f in request.files.getlist('images'):
            if f and f.filename and allowed(f.filename, 'image'):
                existing_imgs.append(save_upload(f, 'solar'))
        mt = request.form.get('mobile_token', '')
        if mt and mt in MOBILE_SESSIONS:
            existing_imgs.extend(MOBILE_SESSIONS.pop(mt)['files'])
        s.images = json.dumps(existing_imgs)

        existing_vids = s.get_videos()
        for f in request.files.getlist('videos'):
            if f and f.filename and allowed(f.filename, 'video'):
                existing_vids.append(save_upload(f, 'solar'))
        s.videos = json.dumps(existing_vids)

        # Keep salesperson sale records in sync with the updated price
        SalespersonSale.query.filter_by(product_type='solar', product_id=s.id).update(
            {'sale_amount': s.price}, synchronize_session=False
        )

        db.session.commit()
        listing_url = f'https://chinacarsinghana.com/solar/{s.id}'
        threading.Thread(
            target=send_new_listing_alert,
            args=('solar', s.name, s.price, listing_url, 'updated'),
            daemon=True
        ).start()
        threading.Thread(
            target=send_price_change_alert,
            args=('solar', s.id, s.name, _old_price, s.price),
            daemon=True
        ).start()
        flash(f'Solar system "{s.name}" updated successfully!', 'success')
        return redirect(url_for('admin_solar'))
    return render_template('admin/solar_form.html', solar=s, action='Edit')


@app.route('/admin/solar/<int:sid>/notify-applicants', methods=['POST'])
@admin_required
def admin_notify_solar_applicants(sid):
    s = SolarSystem.query.get_or_404(sid)
    threading.Thread(
        target=send_manual_price_alert,
        args=('solar', s.id, s.name, s.price),
        daemon=True
    ).start()
    flash(f'Price alert sent to all loan applicants for "{s.name}".', 'success')
    return redirect(url_for('admin_edit_solar', sid=sid))


@app.route('/admin/solar/delete/<int:sid>', methods=['POST'])
@admin_required
def admin_delete_solar(sid):
    s = SolarSystem.query.get_or_404(sid)
    db.session.delete(s)
    db.session.commit()
    flash('Solar system deleted.', 'success')
    return redirect(url_for('admin_solar'))


@app.route('/admin/solar/<int:sid>/delete-media', methods=['POST'])
@admin_required
def admin_delete_solar_media(sid):
    s    = SolarSystem.query.get_or_404(sid)
    path = request.json.get('path')
    kind = request.json.get('kind', 'image')
    if kind == 'image':
        items = [i for i in s.get_images() if i != path]
        s.images = json.dumps(items)
    else:
        items = [i for i in s.get_videos() if i != path]
        s.videos = json.dumps(items)
    db.session.commit()
    _delete_media_file(path)
    return jsonify({'ok': True})


# ── Applications & Contacts ── #

@app.route('/admin/applications')
@admin_required
def admin_applications():
    status = request.args.get('status', '')
    q = LoanApplication.query
    if status: q = q.filter_by(status=status)
    apps = q.order_by(LoanApplication.created_at.desc()).all()
    return render_template('admin/applications.html', applications=apps, current_status=status)


@app.route('/admin/applications/<int:aid>/status', methods=['POST'])
@admin_required
def admin_update_app_status(aid):
    a = LoanApplication.query.get_or_404(aid)
    a.status = request.form.get('status', 'pending')
    db.session.commit()
    flash(f'Application status updated to "{a.status}".', 'success')
    return redirect(url_for('admin_applications'))


@app.route('/admin/applications/<int:aid>')
@admin_required
def admin_view_application(aid):
    a = LoanApplication.query.get_or_404(aid)
    return render_template('admin/application_detail.html', app=a)


@app.route('/admin/applications/<int:aid>/delete', methods=['POST'])
@admin_required
def admin_delete_application(aid):
    a = LoanApplication.query.get_or_404(aid)
    # Unlink any salesperson sale records that reference this application
    SalespersonSale.query.filter_by(loan_application_id=aid).update(
        {'loan_application_id': None}, synchronize_session=False
    )
    db.session.delete(a)
    db.session.commit()
    flash('Loan application deleted.', 'success')
    return redirect(url_for('admin_applications'))


@app.route('/admin/contacts')
@admin_required
def admin_contacts():
    contacts = Contact.query.order_by(Contact.created_at.desc()).all()
    Contact.query.filter_by(status='unread').update({'status': 'read'})
    db.session.commit()
    return render_template('admin/contacts.html', contacts=contacts)


@app.route('/admin/contacts/delete/<int:cid>', methods=['POST'])
@admin_required
def admin_delete_contact(cid):
    c = Contact.query.get_or_404(cid)
    db.session.delete(c)
    db.session.commit()
    flash('Message deleted.', 'success')
    return redirect(url_for('admin_contacts'))


@app.route('/admin/contacts/<int:cid>/reply', methods=['POST'])
@admin_required
def admin_reply_contact(cid):
    c = Contact.query.get_or_404(cid)
    reply = request.form.get('reply', '').strip()
    if not reply:
        flash('Reply cannot be empty.', 'error')
        return redirect(url_for('admin_contacts'))

    c.reply_text = reply
    c.replied_at = datetime.utcnow()
    c.status     = 'replied'
    db.session.commit()

    if SMTP_USER and SMTP_PASSWORD and c.email:
        try:
            subj = f"Re: {c.subject or 'Your Message'} — {BUSINESS['name']}"
            body = (
                f"Hello {c.name},\n\n"
                f"Thank you for contacting {BUSINESS['name']}. "
                f"Here is our reply to your message:\n\n"
                f"{'─' * 50}\n"
                f"Your original message:\n{c.message}\n"
                f"{'─' * 50}\n\n"
                f"{reply}\n\n"
                f"Best regards,\n"
                f"{BUSINESS['name']} Team\n"
                f"📞 {BUSINESS['phone']}\n"
                f"✉️  {BUSINESS['email']}\n"
            )
            msg = MIMEMultipart('alternative')
            msg['Subject'] = subj
            msg['From']    = SMTP_USER
            msg['To']      = c.email
            msg.attach(MIMEText(body, 'plain'))
            with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
                server.ehlo()
                server.starttls()
                server.login(SMTP_USER, SMTP_PASSWORD)
                server.sendmail(SMTP_USER, [c.email], msg.as_string())
            flash(f'Reply sent to {c.name} ({c.email}).', 'success')
        except Exception as e:
            print(f'[EMAIL] Reply send failed: {e}')
            flash('Reply saved but email could not be delivered — check SMTP settings.', 'warning')
    else:
        flash('Reply saved (SMTP not configured — no email sent).', 'success')

    return redirect(url_for('admin_contacts'))


# ── Salesperson Management ── #

@app.route('/admin/salespersons')
@admin_required
def admin_salespersons():
    people = Salesperson.query.order_by(Salesperson.created_at.desc()).all()
    return render_template('admin/salespersons.html', salespersons=people)


@app.route('/admin/salespersons/add', methods=['GET', 'POST'])
@admin_required
def admin_add_salesperson():
    if request.method == 'POST':
        code = request.form.get('code', '').strip().upper()
        if Salesperson.query.filter_by(code=code).first():
            flash('That code is already taken. Choose a unique code.', 'error')
            return redirect(url_for('admin_add_salesperson'))
        sp = Salesperson(
            full_name      = request.form.get('full_name', '').strip(),
            email          = request.form.get('email', '').strip(),
            phone          = request.form.get('phone', '').strip(),
            code           = code,
            commission_pct = request.form.get('commission_pct', type=float) or 1.0,
            notes          = request.form.get('notes', '').strip(),
            active         = 'active' in request.form,
        )
        sp.set_password(request.form.get('password', ''))
        db.session.add(sp)
        db.session.commit()
        if sp.email:
            send_salesperson_welcome(sp)
            flash(f'Salesperson "{sp.full_name}" added (code {sp.code}) — welcome email sent.', 'success')
        else:
            flash(f'Salesperson "{sp.full_name}" added with code {sp.code}.', 'success')
        return redirect(url_for('admin_salespersons'))
    return render_template('admin/salesperson_form.html', sp=None, action='Add')


@app.route('/admin/salespersons/edit/<int:sid>', methods=['GET', 'POST'])
@admin_required
def admin_edit_salesperson(sid):
    sp = Salesperson.query.get_or_404(sid)
    if request.method == 'POST':
        new_code = request.form.get('code', '').strip().upper()
        existing = Salesperson.query.filter_by(code=new_code).first()
        if existing and existing.id != sp.id:
            flash('That code is already taken.', 'error')
            return redirect(url_for('admin_edit_salesperson', sid=sid))
        sp.full_name      = request.form.get('full_name', '').strip()
        sp.email          = request.form.get('email', '').strip()
        sp.phone          = request.form.get('phone', '').strip()
        sp.code           = new_code
        sp.commission_pct = request.form.get('commission_pct', type=float) or 1.0
        sp.notes          = request.form.get('notes', '').strip()
        sp.active         = 'active' in request.form
        new_pw = request.form.get('new_password', '').strip()
        if new_pw:
            sp.set_password(new_pw)
        db.session.commit()
        flash('Salesperson updated.', 'success')
        return redirect(url_for('admin_salespersons'))
    return render_template('admin/salesperson_form.html', sp=sp, action='Edit')


@app.route('/admin/salespersons/delete/<int:sid>', methods=['POST'])
@admin_required
def admin_delete_salesperson(sid):
    sp = Salesperson.query.get_or_404(sid)
    db.session.delete(sp)
    db.session.commit()
    flash('Salesperson removed.', 'success')
    return redirect(url_for('admin_salespersons'))


@app.route('/admin/salespersons/<int:sid>/sales')
@admin_required
def admin_salesperson_sales(sid):
    sp    = Salesperson.query.get_or_404(sid)
    sales = SalespersonSale.query.filter_by(salesperson_id=sid).order_by(SalespersonSale.created_at.desc()).all()
    return render_template('admin/salesperson_sales.html', sp=sp, sales=sales)


@app.route('/admin/salespersons/<int:sid>/sales/add', methods=['GET', 'POST'])
@admin_required
def admin_add_sale(sid):
    sp = Salesperson.query.get_or_404(sid)
    if request.method == 'POST':
        sale = SalespersonSale(
            salesperson_id   = sid,
            product_type     = request.form.get('product_type'),
            product_id       = request.form.get('product_id', type=int),
            product_name     = request.form.get('product_name', '').strip(),
            sale_amount      = request.form.get('sale_amount', type=float),
            customer_name    = request.form.get('customer_name', '').strip(),
            customer_email   = request.form.get('customer_email', '').strip(),
            customer_phone   = request.form.get('customer_phone', '').strip(),
            notes            = request.form.get('notes', '').strip(),
            status           = request.form.get('status', 'pending'),
        )
        db.session.add(sale)
        db.session.commit()
        flash('Sale recorded.', 'success')
        return redirect(url_for('admin_salesperson_sales', sid=sid))
    vehicles    = Vehicle.query.filter_by(available=True).all()
    solar_list  = SolarSystem.query.filter_by(available=True).all()
    return render_template('admin/add_sale.html', sp=sp, vehicles=vehicles, solar_systems=solar_list)


@app.route('/admin/sales/<int:sale_id>/status', methods=['POST'])
@admin_required
def admin_update_sale_status(sale_id):
    sale = SalespersonSale.query.get_or_404(sale_id)
    sale.status = request.form.get('status', 'pending')
    db.session.commit()
    flash('Sale status updated.', 'success')
    return redirect(url_for('admin_salesperson_sales', sid=sale.salesperson_id))


# ── Salesperson self-registration ── #

@app.route('/sales/register', methods=['GET', 'POST'])
@app.route('/sales/register/<token>', methods=['GET', 'POST'])
def sales_register(token=None):
    """Public self-registration page — open signup or token-gated via admin link."""
    # If a token was supplied, validate it; if no token, allow open registration
    if token:
        try:
            setting = SiteSettings.query.filter_by(key='sp_reg_token').first()
        except Exception:
            setting = None
        if not setting or setting.value != token:
            return '<h2 style="font-family:sans-serif;text-align:center;padding:3rem;color:#dc2626;">This registration link is no longer valid. Please contact the admin for a new one.</h2>', 403

    if request.method == 'POST':
        full_name = request.form.get('full_name', '').strip()
        email     = request.form.get('email', '').strip()
        phone     = request.form.get('phone', '').strip()
        password  = request.form.get('password', '')
        confirm   = request.form.get('confirm', '')

        if not full_name or not password:
            flash('Name and password are required.', 'error')
        elif password != confirm:
            flash('Passwords do not match.', 'error')
        elif len(password) < 6:
            flash('Password must be at least 6 characters.', 'error')
        elif Salesperson.query.filter(
                db.func.lower(Salesperson.full_name) == full_name.lower()).first():
            flash('An account with this name already exists. If this is you, please log in or contact the admin.', 'error')
        elif email and Salesperson.query.filter(
                db.func.lower(Salesperson.email) == email.lower()).first():
            flash('An account with this email address already exists. Please contact admin if you need help.', 'error')
        elif phone and Salesperson.query.filter_by(phone=phone).first():
            flash('An account with this phone number already exists. Please contact admin if you need help.', 'error')
        else:
            # Auto-generate unique code from initials + random digits
            base = ''.join(w[0] for w in full_name.split() if w).upper()[:3]
            while True:
                code = base + str(uuid.uuid4().int)[:4]
                if not Salesperson.query.filter_by(code=code).first():
                    break
            sp = Salesperson(
                full_name      = full_name,
                email          = email,
                phone          = phone,
                code           = code,
                active         = True,
                commission_pct = 1.0,
            )
            sp.set_password(password)
            db.session.add(sp)
            db.session.commit()
            # Email the salesperson their code + admin notification
            if email:
                send_salesperson_welcome(sp)
                flash(f'Registration successful! Your code is {code}. Check your email for login details.', 'success')
            else:
                flash(f'Registration successful! Your code is {code}. Please log in.', 'success')
            return redirect(url_for('sales_login'))

    return render_template('sales/register.html')


@app.route('/admin/salespersons/reg-link', methods=['POST'])
@admin_required
def admin_gen_reg_link():
    """Generate a new salesperson registration link (valid 7 days)."""
    token = uuid.uuid4().hex
    row   = SiteSettings.query.filter_by(key='sp_reg_token').first()
    if row:
        row.value = token
    else:
        db.session.add(SiteSettings(key='sp_reg_token', value=token))
    db.session.commit()
    link = f"https://chinacarsinghana.com/sales/register/{token}"
    return jsonify({'link': link})


@app.route('/admin/team/email', methods=['GET', 'POST'])
@admin_required
def admin_team_email():
    """Compose and send a custom email to selected salespersons."""
    salespersons = Salesperson.query.filter_by(active=True).order_by(Salesperson.full_name).all()

    if request.method == 'POST':
        subject    = request.form.get('subject', '').strip()
        body_tpl   = request.form.get('body', '').strip()
        recipient_ids = request.form.getlist('recipients')

        if not subject or not body_tpl:
            flash('Subject and message body are required.', 'error')
            return render_template('admin/team_email.html', salespersons=salespersons)

        if not recipient_ids:
            flash('Please select at least one recipient.', 'error')
            return render_template('admin/team_email.html', salespersons=salespersons)

        if not SMTP_USER or not SMTP_PASSWORD:
            flash('SMTP is not configured — emails cannot be sent.', 'error')
            return render_template('admin/team_email.html', salespersons=salespersons)

        targets = [sp for sp in salespersons if str(sp.id) in recipient_ids and sp.email]
        no_email = [sp for sp in salespersons if str(sp.id) in recipient_ids and not sp.email]

        sent = 0
        failed = 0
        for sp in targets:
            # Replace merge tags
            personalised_body = (body_tpl
                .replace('{name}', sp.full_name)
                .replace('{code}', sp.code)
                .replace('{ref_link}', f'https://chinacarsinghana.com?ref={sp.code}')
                .replace('{login_url}', 'https://chinacarsinghana.com/sales/login')
                .replace('{commission}', f'{sp.commission_pct}%'))

            personalised_body += (
                f"\n\n{'─' * 50}"
                f"\nChina Cars in Ghana — Sales Team"
                f"\nhttps://chinacarsinghana.com/sales/login"
            )

            msg = MIMEMultipart('alternative')
            msg['Subject'] = subject
            msg['From']    = SMTP_USER
            msg['To']      = sp.email
            msg.attach(MIMEText(personalised_body, 'plain'))
            try:
                with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
                    server.ehlo(); server.starttls()
                    server.login(SMTP_USER, SMTP_PASSWORD)
                    server.sendmail(SMTP_USER, [sp.email], msg.as_string())
                sent += 1
            except Exception as e:
                print(f'[EMAIL] Team email to {sp.email} failed: {e}')
                failed += 1

        parts = []
        if sent:
            parts.append(f'Email sent to {sent} salesperson{"s" if sent != 1 else ""}.')
        if failed:
            parts.append(f'{failed} failed (check server logs).')
        if no_email:
            names = ', '.join(sp.full_name for sp in no_email)
            parts.append(f'Skipped (no email on file): {names}.')
        flash(' '.join(parts), 'success' if sent else 'error')
        return redirect(url_for('admin_team_email'))

    return render_template('admin/team_email.html', salespersons=salespersons)


def _build_proposal_pptx(sp_name='[YOUR FULL NAME]', sp_code='[YOUR CODE]',
                          sp_phone='[YOUR PHONE NUMBER]', sp_email='[YOUR EMAIL ADDRESS]'):
    """Generate the China Cars in Ghana business proposal as a PPTX in memory."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    NAVY   = RGBColor(0x0F, 0x17, 0x2A)
    GOLD   = RGBColor(0xF5, 0x9E, 0x0B)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    LGRAY  = RGBColor(0xF1, 0xF5, 0xF9)
    BLUE   = RGBColor(0x1D, 0x4E, 0xD8)
    GREEN  = RGBColor(0x16, 0xA3, 0x4A)
    DKGRAY = RGBColor(0x1E, 0x29, 0x3B)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # fully blank

    def px(val): return Inches(val)

    def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
        shape = slide.shapes.add_shape(1, px(l), px(t), px(w), px(h))
        shape.line.fill.background()
        if fill_rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
        else:
            shape.fill.background()
        if line_rgb:
            from pptx.util import Pt as Pt2
            shape.line.color.rgb = line_rgb
            shape.line.width = Pt2(line_width or 1)
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, text, l, t, w, h,
                 font_size=16, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
        txb = slide.shapes.add_textbox(px(l), px(t), px(w), px(h))
        txb.word_wrap = wrap
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = 'Calibri'
        return txb

    def add_multiline(slide, lines, l, t, w, h,
                      font_size=13, bold_first=False, color=WHITE,
                      align=PP_ALIGN.LEFT, line_spacing=1.1):
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        txb = slide.shapes.add_textbox(px(l), px(t), px(w), px(h))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.size = Pt2(font_size)
            run.font.bold = (bold_first and i == 0)
            run.font.color.rgb = color
            run.font.name = 'Calibri'
        return txb

    # ── SLIDE 1: COVER ────────────────────────────────────────────────────── #
    s1 = prs.slides.add_slide(blank_layout)

    # Full navy background
    add_rect(s1, 0, 0, 13.33, 7.5, fill_rgb=NAVY)

    # Gold accent bar left edge
    add_rect(s1, 0, 0, 0.18, 7.5, fill_rgb=GOLD)

    # Gold accent bar bottom
    add_rect(s1, 0, 6.9, 13.33, 0.6, fill_rgb=GOLD)

    # Dark panel right side
    add_rect(s1, 8.5, 0, 4.83, 7.5, fill_rgb=DKGRAY)

    # Top-left: "PARTNERSHIP PROPOSAL"
    add_text(s1, 'BUSINESS PARTNERSHIP PROPOSAL',
             0.45, 0.35, 7.8, 0.55,
             font_size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Main headline
    add_text(s1, 'China Cars\nin Ghana',
             0.45, 0.85, 7.5, 2.2,
             font_size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Tag line
    add_text(s1, 'Premium Vehicles · Solar Energy Systems · Flexible Finance',
             0.45, 3.15, 7.8, 0.5,
             font_size=13, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

    # Divider line under tag
    add_rect(s1, 0.45, 3.72, 4.5, 0.04, fill_rgb=GOLD)

    # Prepared for block
    add_text(s1, 'PREPARED FOR:',
             0.45, 3.9, 4, 0.35,
             font_size=9, bold=True, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)
    add_text(s1, '[COMPANY NAME]',
             0.45, 4.2, 6, 0.6,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Date
    add_text(s1, datetime.utcnow().strftime('%B %Y'),
             0.45, 4.85, 4, 0.35,
             font_size=10, bold=False, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)

    # Website bottom bar
    add_text(s1, 'www.chinacarsinghana.com',
             0.45, 7.0, 5, 0.35,
             font_size=10, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # Right panel content
    add_text(s1, '⭐ WHY US',
             8.75, 0.6, 4.2, 0.4,
             font_size=10, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    right_points = [
        '✔  Wide range of premium Chinese vehicles',
        '✔  Solar energy systems for homes & business',
        '✔  Republic Bank financing — easy approval',
        '✔  Flexible payment plans available',
        '✔  Exclusive discounts for referred customers',
        '✔  Trusted dealer in Accra, Ghana',
        '✔  Dedicated sales support team',
        '✔  After-sales service & warranty',
    ]
    add_multiline(s1, right_points,
                  8.75, 1.1, 4.2, 4.5,
                  font_size=11, color=LGRAY, align=PP_ALIGN.LEFT)

    add_text(s1, 'Your Trusted Partner in Ghana',
             8.75, 6.3, 4.2, 0.45,
             font_size=10, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # ── SLIDE 2: SERVICES & LOAN PACKAGE ──────────────────────────────────── #
    s2 = prs.slides.add_slide(blank_layout)

    # Background
    add_rect(s2, 0, 0, 13.33, 7.5, fill_rgb=LGRAY)

    # Header band
    add_rect(s2, 0, 0, 13.33, 1.25, fill_rgb=NAVY)
    add_rect(s2, 0, 1.25, 13.33, 0.06, fill_rgb=GOLD)

    add_text(s2, 'Our Products & Services',
             0.4, 0.25, 9, 0.7,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s2, 'SLIDE 2 OF 3  |  chinacarsinghana.com',
             10.0, 0.45, 3.0, 0.4,
             font_size=8, bold=False, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.RIGHT)

    # ── Card 1: Vehicles ──
    add_rect(s2, 0.35, 1.5, 3.85, 5.5, fill_rgb=WHITE)
    add_rect(s2, 0.35, 1.5, 3.85, 0.45, fill_rgb=BLUE)
    add_text(s2, '🚗  PREMIUM VEHICLES',
             0.45, 1.55, 3.65, 0.35,
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    v_lines = [
        'CATEGORIES AVAILABLE:',
        '  Saloon / Sedan Cars',
        '  SUVs & Crossovers',
        '  Pick-Up Trucks & 4x4s',
        '  Minivans & MPVs',
        '  Buses & Minibuses',
        '  Farm & Off-Road Vehicles',
        '  Electric Vehicles (EV)',
        '  Petrol & Diesel Options',
        '',
        '• Brand-new, factory-direct pricing',
        '• Full Ghana documentation & customs clearance',
        '• Test drives available on request',
        '• Corporate fleet & bulk order pricing',
        '• Finance via Republic Bank available',
    ]
    add_multiline(s2, v_lines, 0.45, 2.1, 3.65, 4.8,
                  font_size=9.8, color=DKGRAY, align=PP_ALIGN.LEFT)

    # ── Card 2: Solar ──
    add_rect(s2, 4.55, 1.5, 3.85, 5.5, fill_rgb=WHITE)
    add_rect(s2, 4.55, 1.5, 3.85, 0.45, fill_rgb=GOLD)
    add_text(s2, '☀️  SOLAR ENERGY SYSTEMS',
             4.65, 1.55, 3.65, 0.35,
             font_size=10, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    s_lines = [
        'Residential, Commercial & Industrial',
        '',
        '• Complete solar panel installations',
        '• Battery backup & inverter systems',
        '• High-efficiency panels & equipment',
        '• Energy independence from grid outages',
        '• Reduce electricity bills significantly',
        '• Ideal for offices, factories & estates',
        '• Maintenance & warranty support',
    ]
    add_multiline(s2, s_lines, 4.65, 2.1, 3.65, 4.6,
                  font_size=10.5, color=DKGRAY, align=PP_ALIGN.LEFT)

    # ── Card 3: Finance ──
    add_rect(s2, 8.75, 1.5, 4.2, 5.5, fill_rgb=WHITE)
    add_rect(s2, 8.75, 1.5, 4.2, 0.45, fill_rgb=GREEN)
    add_text(s2, '🏦  REPUBLIC BANK FINANCE',
             8.85, 1.55, 4.0, 0.35,
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    f_lines = [
        'Flexible Loan Packages',
        '',
        '• Finance up to 100% of vehicle/system cost',
        '• Loan terms: 12 – 84 months',
        '• Competitive interest rates',
        '• Fast approval — 24 to 48 hours',
        '• Minimal documentation required',
        '• Available to employed & self-employed',
        '• Joint applications accepted',
        '• Exclusive discount for referred customers',
        '',
        'Apply online at:',
        'chinacarsinghana.com/loan',
    ]
    add_multiline(s2, f_lines, 8.85, 2.1, 3.95, 4.6,
                  font_size=10.5, color=DKGRAY, align=PP_ALIGN.LEFT)

    # ── SLIDE 3: CONTACT & NEXT STEPS ─────────────────────────────────────── #
    s3 = prs.slides.add_slide(blank_layout)

    # Navy background full
    add_rect(s3, 0, 0, 13.33, 7.5, fill_rgb=NAVY)

    # Gold side bar
    add_rect(s3, 0, 0, 0.18, 7.5, fill_rgb=GOLD)

    # Gold bottom bar
    add_rect(s3, 0, 6.9, 13.33, 0.6, fill_rgb=GOLD)

    # Header
    add_text(s3, "LET'S WORK TOGETHER",
             0.45, 0.3, 9, 0.45,
             font_size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_text(s3, 'Get in Touch',
             0.45, 0.72, 9, 0.85,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(s3, 0.45, 1.6, 4.0, 0.05, fill_rgb=GOLD)

    # Left: Salesperson contact block
    add_text(s3, 'YOUR DEDICATED SALES REPRESENTATIVE',
             0.45, 1.8, 6.2, 0.35,
             font_size=9, bold=True, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)

    sp_fields = [
        ('Name',          sp_name),
        ('Phone',         sp_phone),
        ('Email',         sp_email),
        ('Sales Code',    sp_code),
        ('Website',       'www.chinacarsinghana.com'),
        ('Email Enquiry', 'kofi@chinacarsinghana.com'),
    ]
    y = 2.2
    for label, val in sp_fields:
        add_text(s3, label.upper(),
                 0.45, y, 2.0, 0.28,
                 font_size=8, bold=True, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)
        add_text(s3, val,
                 2.5, y, 4.5, 0.32,
                 font_size=11, bold=(val.startswith('[')), color=GOLD if val.startswith('[') else WHITE,
                 align=PP_ALIGN.LEFT)
        y += 0.45

    # Right: Next steps
    add_rect(s3, 7.2, 1.7, 5.7, 4.8, fill_rgb=DKGRAY)
    add_text(s3, 'NEXT STEPS',
             7.45, 1.95, 5.2, 0.4,
             font_size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    steps = [
        '01  Schedule a product demonstration',
        '02  Explore our full vehicle & solar inventory',
        '03  Receive a tailored corporate quote',
        '04  Apply for Republic Bank financing',
        '05  Place your order — delivery in days',
    ]
    add_multiline(s3, steps, 7.45, 2.45, 5.1, 3.5,
                  font_size=11.5, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(s3, 'Visit us online to browse inventory,\napply for finance & get instant quotes.',
             7.45, 5.5, 5.1, 0.7,
             font_size=10, bold=False, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)

    # Bottom bar text
    add_text(s3, 'chinacarsinghana.com  |  Accra, Ghana  |  Trusted Dealer Since Establishment',
             0.45, 7.0, 9, 0.35,
             font_size=9, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@app.route('/admin/team/proposal', methods=['GET', 'POST'])
@admin_required
def admin_team_proposal():
    """Send the business proposal PPTX to selected salespersons."""
    salespersons = Salesperson.query.filter_by(active=True).order_by(Salesperson.full_name).all()

    if request.method == 'POST':
        recipient_ids = request.form.getlist('recipients')

        if not recipient_ids:
            flash('Please select at least one recipient.', 'error')
            return render_template('admin/team_proposal.html', salespersons=salespersons)

        if not SMTP_USER or not SMTP_PASSWORD:
            flash('SMTP is not configured — emails cannot be sent.', 'error')
            return render_template('admin/team_proposal.html', salespersons=salespersons)

        targets  = [sp for sp in salespersons if str(sp.id) in recipient_ids and sp.email]
        no_email = [sp for sp in salespersons if str(sp.id) in recipient_ids and not sp.email]

        try:
            pptx_buf = _build_proposal_pptx()
        except Exception as e:
            flash(f'Failed to generate presentation: {e}', 'error')
            return render_template('admin/team_proposal.html', salespersons=salespersons)

        target_companies = {
            'Banks & Financial Institutions': [
                'GCB Bank', 'Ecobank Ghana', 'Fidelity Bank Ghana', 'Absa Bank Ghana',
                'Standard Chartered Ghana', 'Republic Bank Ghana', 'UBA Ghana',
                'Zenith Bank Ghana', 'Access Bank Ghana', 'Consolidated Bank Ghana',
                'First Atlantic Bank', 'OmniBank Ghana', 'CalBank',
            ],
            'Telecoms & Technology': [
                'MTN Ghana', 'Vodafone Ghana', 'AirtelTigo', 'Huawei Ghana',
            ],
            'Government Agencies & Civil Service': [
                'Ghana Revenue Authority', 'SSNIT', 'ECG (Electricity Company of Ghana)',
                'Ghana Water Company', 'Ghana Ports & Harbours Authority',
                'Ghana Highways Authority', 'Accra Metropolitan Assembly',
                'Ghana Health Service', 'Ghana Police Service', 'Ghana Armed Forces',
                'Ghana Immigration Service', 'Passport Office', 'DVLA Ghana',
                'Ministry of Roads & Highways', 'Ministry of Finance',
                'Lands Commission', 'Ghana National Fire Service',
            ],
            'Energy & Mining': [
                'Tullow Oil Ghana', 'AngloGold Ashanti', 'Newmont Ghana',
                'Kosmos Energy', 'Total Energies Ghana', 'GNPC',
            ],
            'Education': [
                'University of Ghana', 'KNUST', 'Ashesi University',
                'University of Cape Coast', 'University of Professional Studies',
            ],
            'Healthcare': [
                'Korle Bu Teaching Hospital', '37 Military Hospital',
                'Ridge Hospital', 'Catholic Health Service',
            ],
            'FMCG & Corporates': [
                'Unilever Ghana', 'Nestlé Ghana', 'Fan Milk (Danone)',
                'PZ Cussons', 'Guinness Ghana', 'Olam Ghana', 'Kasapreko',
            ],
        }
        companies_text = ''
        for category, names in target_companies.items():
            companies_text += f'\n[{category}]\n'
            companies_text += '\n'.join(f'  • {c}' for c in names)
            companies_text += '\n'

        sent = 0
        failed = 0
        pptx_buf.seek(0)
        pptx_bytes = pptx_buf.read()

        for sp in targets:
            email_body = f"""\
Hi {sp.full_name},

Please find attached a business proposal presentation (PowerPoint) for China Cars in Ghana.

HOW TO USE THIS PROPOSAL
─────────────────────────
1. Open the attached file in Microsoft PowerPoint or Google Slides.
2. On the COVER slide, replace [COMPANY NAME] with the name of the company you are pitching to.
3. On the CONTACT slide (Slide 3), fill in your:
     • Full Name
     • Phone Number
     • Email Address
     • Your Sales Code: {sp.code}
4. Save the file and send or present it professionally.

SUGGESTED TARGET COMPANIES IN GHANA
─────────────────────────────────────
The following organisations are excellent prospects — they have staff, fleets, and energy needs:

{companies_text}

Focus on their procurement, fleet management, facilities, or HR departments. A well-placed introduction can lead to bulk vehicle orders, corporate solar installations, or group loan applications through Republic Bank.

YOUR SALES DASHBOARD
─────────────────────
Log in to your dashboard for the full target company list, Ghana sales tips,
and to download this proposal again at any time:

  🖥  https://chinacarsinghana.com/sales/portal

Your dashboard also lets you:
  • Browse current vehicle & solar inventory
  • Submit new customer leads
  • Track your sales and commission in real time
  • Re-download this proposal whenever you need it

NEED HELP OR HAVE QUESTIONS?
──────────────────────────────
Contact the admin team directly at:
  📧  kofi@chinacarsinghana.com

We are happy to assist you with custom quotes, additional materials, or to join a client meeting.

Your referral link:
  🔗  https://chinacarsinghana.com?ref={sp.code}

Good luck — let's close some big deals!
— China Cars in Ghana Admin Team
{'─' * 54}
www.chinacarsinghana.com · Accra, Ghana
"""
            msg = MIMEMultipart()
            msg['Subject'] = 'Your Business Proposal — China Cars in Ghana [Ready to Send]'
            msg['From']    = SMTP_USER
            msg['To']      = sp.email
            msg.attach(MIMEText(email_body, 'plain'))

            part = MIMEBase('application',
                            'vnd.openxmlformats-officedocument.presentationml.presentation')
            part.set_payload(pptx_bytes)
            encoders.encode_base64(part)
            part.add_header('Content-Disposition',
                            'attachment',
                            filename='ChinaCarsGhana_Proposal.pptx')
            msg.attach(part)

            try:
                with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
                    server.ehlo(); server.starttls()
                    server.login(SMTP_USER, SMTP_PASSWORD)
                    server.sendmail(SMTP_USER, [sp.email], msg.as_string())
                sent += 1
            except Exception as e:
                print(f'[EMAIL] Proposal to {sp.email} failed: {e}')
                failed += 1

        parts = []
        if sent:
            parts.append(f'Proposal sent to {sent} salesperson{"s" if sent != 1 else ""}.')
        if failed:
            parts.append(f'{failed} failed (check server logs).')
        if no_email:
            names = ', '.join(sp.full_name for sp in no_email)
            parts.append(f'Skipped (no email on file): {names}.')
        flash(' '.join(parts), 'success' if sent else 'error')
        return redirect(url_for('admin_team_proposal'))

    return render_template('admin/team_proposal.html', salespersons=salespersons)


# ── Salesperson Portal (their own login) ── #

@app.route('/sales/login', methods=['GET', 'POST'])
def sales_login():
    if 'sales_id' in session:
        return redirect(url_for('sales_portal'))
    if request.method == 'POST':
        code = request.form.get('code', '').strip().upper()
        pw   = request.form.get('password', '')
        sp   = Salesperson.query.filter_by(code=code, active=True).first()
        if sp and sp.check_password(pw):
            session['sales_id']   = sp.id
            session['sales_name'] = sp.full_name
            session['sales_code'] = sp.code
            return redirect(url_for('sales_portal'))
        flash('Invalid code or password.', 'error')
    return render_template('sales/login.html')


@app.route('/sales/forgot-password', methods=['GET', 'POST'])
def sales_forgot_password():
    if request.method == 'POST':
        code  = request.form.get('code', '').strip().upper()
        email = request.form.get('email', '').strip().lower()
        sp = Salesperson.query.filter_by(code=code, active=True).first()
        if sp and sp.email and sp.email.lower() == email:
            token = uuid.uuid4().hex
            sp.reset_token        = token
            sp.reset_token_expiry = datetime.utcnow() + timedelta(hours=1)
            db.session.commit()
            reset_url = url_for('sales_reset_password', token=token, _external=True)
            if SMTP_USER and SMTP_PASSWORD:
                try:
                    body = f"""\
Hi {sp.full_name},

A password reset was requested for your sales portal account.

Click the link below to set a new password (valid for 1 hour):

  {reset_url}

If you did not request this, ignore this email — your password has not changed.

──────────────────────────────────────
China Cars in Ghana — Sales Team
"""
                    msg = MIMEMultipart('alternative')
                    msg['Subject'] = 'Reset Your Sales Portal Password'
                    msg['From']    = SMTP_USER
                    msg['To']      = sp.email
                    msg.attach(MIMEText(body, 'plain'))
                    def _send():
                        try:
                            with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
                                server.ehlo(); server.starttls()
                                server.login(SMTP_USER, SMTP_PASSWORD)
                                server.sendmail(SMTP_USER, [sp.email], msg.as_string())
                        except Exception as e:
                            print(f'[EMAIL] Password reset email failed: {e}')
                    threading.Thread(target=_send, daemon=True).start()
                except Exception as e:
                    print(f'[EMAIL] Password reset prepare failed: {e}')
        # Always show success to prevent email enumeration
        flash('If that code and email match our records, a reset link has been sent. Check your inbox.', 'success')
        return redirect(url_for('sales_forgot_password'))
    return render_template('sales/forgot_password.html')


@app.route('/sales/reset-password/<token>', methods=['GET', 'POST'])
def sales_reset_password(token):
    sp = Salesperson.query.filter_by(reset_token=token).first()
    if not sp or not sp.reset_token_expiry or sp.reset_token_expiry < datetime.utcnow():
        flash('This reset link is invalid or has expired. Please request a new one.', 'error')
        return redirect(url_for('sales_forgot_password'))
    if request.method == 'POST':
        pw  = request.form.get('password', '')
        pw2 = request.form.get('password2', '')
        if len(pw) < 6:
            flash('Password must be at least 6 characters.', 'error')
            return render_template('sales/reset_password.html', token=token)
        if pw != pw2:
            flash('Passwords do not match.', 'error')
            return render_template('sales/reset_password.html', token=token)
        sp.set_password(pw)
        sp.reset_token        = None
        sp.reset_token_expiry = None
        db.session.commit()
        flash('Password updated! You can now log in with your new password.', 'success')
        return redirect(url_for('sales_login'))
    return render_template('sales/reset_password.html', token=token)


@app.route('/sales/logout')
def sales_logout():
    session.pop('sales_id', None)
    session.pop('sales_name', None)
    session.pop('sales_code', None)
    return redirect(url_for('sales_login'))


def sales_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if 'sales_id' not in session:
            return redirect(url_for('sales_login'))
        return f(*args, **kwargs)
    return decorated


@app.route('/sales')
@app.route('/sales/portal')
@sales_required
def sales_portal():
    sp = Salesperson.query.get(session['sales_id'])
    if not sp:
        session.clear()
        return redirect(url_for('sales_login'))
    sales = SalespersonSale.query.filter_by(salesperson_id=sp.id).order_by(SalespersonSale.created_at.desc()).limit(10).all()
    return render_template('sales/portal.html', sp=sp, recent_sales=sales)


@app.route('/sales/proposal/download')
@sales_required
def sales_download_proposal():
    sp = Salesperson.query.get_or_404(session['sales_id'])
    buf = _build_proposal_pptx(
        sp_name  = sp.full_name,
        sp_code  = sp.code,
        sp_phone = sp.phone  or '[YOUR PHONE NUMBER]',
        sp_email = sp.email  or '[YOUR EMAIL ADDRESS]',
    )
    return send_file(
        buf,
        as_attachment=True,
        download_name='ChinaCarsGhana_Proposal.pptx',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )


@app.route('/sales/inventory')
@sales_required
def sales_inventory():
    vehicles    = Vehicle.query.filter_by(available=True).all()
    solar_list  = SolarSystem.query.filter_by(available=True).all()
    return render_template('sales/inventory.html', vehicles=vehicles, solar_systems=solar_list)


@app.route('/sales/submit', methods=['GET', 'POST'])
@sales_required
def sales_submit_lead():
    sp = Salesperson.query.get_or_404(session['sales_id'])
    vehicles   = Vehicle.query.filter_by(available=True).all()
    solar_list = SolarSystem.query.filter_by(available=True).all()
    if request.method == 'POST':
        product_type = request.form.get('product_type')
        product_id   = request.form.get('product_id', type=int)
        product_name = request.form.get('product_name', '').strip()
        sale_amount  = request.form.get('sale_amount', type=float)

        # Also create a loan application if customer needs financing
        if request.form.get('needs_loan') == '1':
            loan_app = LoanApplication(
                first_name        = request.form.get('customer_name', '').split()[0] if request.form.get('customer_name') else '',
                last_name         = ' '.join(request.form.get('customer_name', '').split()[1:]) or '',
                email             = request.form.get('customer_email', '').strip(),
                phone             = request.form.get('customer_phone', '').strip(),
                product_type      = product_type,
                product_id        = product_id,
                product_name      = product_name,
                product_price     = sale_amount,
                loan_amount       = sale_amount,
                message           = f'Lead submitted by salesperson {sp.full_name} (Code: {sp.code})',
            )
            db.session.add(loan_app)
            db.session.flush()
            loan_id = loan_app.id
        else:
            loan_id = None

        sale = SalespersonSale(
            salesperson_id      = sp.id,
            product_type        = product_type,
            product_id          = product_id,
            product_name        = product_name,
            sale_amount         = sale_amount,
            customer_name       = request.form.get('customer_name', '').strip(),
            customer_email      = request.form.get('customer_email', '').strip(),
            customer_phone      = request.form.get('customer_phone', '').strip(),
            loan_application_id = loan_id,
            notes               = request.form.get('notes', '').strip(),
            status              = 'pending',
        )
        db.session.add(sale)
        db.session.commit()
        flash('Lead submitted successfully! The admin will review it shortly.', 'success')
        return redirect(url_for('sales_portal'))
    return render_template('sales/submit_lead.html', sp=sp, vehicles=vehicles, solar_systems=solar_list)


@app.route('/admin/settings', methods=['GET', 'POST'])
@admin_required
def admin_settings():
    admin = AdminUser.query.get(session['admin_id'])
    if request.method == 'POST':
        admin.full_name = request.form.get('full_name', '').strip()
        admin.email     = request.form.get('email', '').strip()
        new_pw          = request.form.get('new_password', '').strip()
        if new_pw:
            if admin.check_password(request.form.get('current_password', '')):
                admin.set_password(new_pw)
                flash('Password updated.', 'success')
            else:
                flash('Current password is incorrect.', 'error')
                return redirect(url_for('admin_settings'))
        # Save GHS rate
        rate_str = request.form.get('ghs_rate', '').strip()
        try:
            rate_val = float(rate_str)
            if rate_val > 0:
                row = SiteSettings.query.filter_by(key='usd_to_ghs_rate').first()
                if row:
                    row.value = str(rate_val)
                else:
                    db.session.add(SiteSettings(key='usd_to_ghs_rate', value=str(rate_val)))
        except (ValueError, TypeError):
            pass
        # Save loan interest rate
        loan_rate_str = request.form.get('loan_interest_rate', '').strip()
        try:
            loan_rate_val = float(loan_rate_str)
            if 0 < loan_rate_val <= 100:
                row = SiteSettings.query.filter_by(key='loan_interest_rate').first()
                if row:
                    row.value = str(loan_rate_val)
                else:
                    db.session.add(SiteSettings(key='loan_interest_rate', value=str(loan_rate_val)))
        except (ValueError, TypeError):
            pass
        db.session.commit()
        session['admin_name'] = admin.full_name or admin.username
        flash('Settings saved successfully.', 'success')
    current_rate      = get_ghs_rate()
    loan_rate_row     = SiteSettings.query.filter_by(key='loan_interest_rate').first()
    current_loan_rate = float(loan_rate_row.value) if loan_rate_row else 7.0
    return render_template('admin/settings.html', admin=admin,
                           current_rate=current_rate,
                           current_loan_rate=current_loan_rate)


# ─────────────────── SOLAR QUOTE SHORT LINKS ────────────────────────────── #

import random, string as _string

def _gen_quote_code():
    chars = _string.ascii_uppercase + _string.digits
    while True:
        code = ''.join(random.choices(chars, k=6))
        if not SolarQuote.query.filter_by(code=code).first():
            return code

@app.route('/api/solar-quote', methods=['POST'])
def api_create_solar_quote():
    try:
        data = request.get_json(force=True)
        if not data or not data.get('a'):
            return jsonify({'error': 'no appliances'}), 400
        code  = _gen_quote_code()
        quote = SolarQuote(code=code, data=json.dumps(data))
        db.session.add(quote)
        db.session.commit()
        short_url = request.host_url.rstrip('/') + '/q/' + code
        return jsonify({'code': code, 'url': short_url})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/q/<code>')
def solar_quote_redirect(code):
    quote = SolarQuote.query.filter_by(code=code).first()
    if not quote:
        return redirect(url_for('index'))
    # Pass code to index so JS can fetch the data via /api/solar-quote/<code>
    featured_vehicles = Vehicle.query.filter_by(featured=True, available=True).limit(4).all()
    featured_solar    = SolarSystem.query.filter_by(featured=True, available=True).limit(2).all()
    latest_vehicles   = Vehicle.query.filter_by(available=True).order_by(Vehicle.created_at.desc()).limit(6).all()
    latest_solar      = SolarSystem.query.filter_by(available=True).order_by(SolarSystem.created_at.desc()).limit(4).all()
    newest_vehicle    = Vehicle.query.filter_by(available=True).order_by(Vehicle.created_at.desc()).first()
    return render_template('index.html',
        featured_vehicles=featured_vehicles,
        featured_solar=featured_solar,
        latest_vehicles=latest_vehicles,
        latest_solar=latest_solar,
        newest_vehicle=newest_vehicle,
        solar_quote_code=code)

@app.route('/api/solar-quote/<code>')
def api_get_solar_quote(code):
    quote = SolarQuote.query.filter_by(code=code).first()
    if not quote:
        return jsonify({'error': 'not found'}), 404
    return jsonify(json.loads(quote.data))


# ─────────────────────────── INIT ───────────────────────────────────────── #

def init_db():
    try:
        with app.app_context():
            db.create_all()

            # Add new columns to existing tables (safe — IF NOT EXISTS skips if already present)
            migrations = [
                "ALTER TABLE loan_application ADD COLUMN IF NOT EXISTS sp_code  VARCHAR(20)",
                "ALTER TABLE loan_application ADD COLUMN IF NOT EXISTS sp_name  VARCHAR(100)",
                "ALTER TABLE loan_application ADD COLUMN IF NOT EXISTS sp_phone VARCHAR(20)",
                "ALTER TABLE car_order        ADD COLUMN IF NOT EXISTS sp_code  VARCHAR(20)",
                "ALTER TABLE car_order        ADD COLUMN IF NOT EXISTS sp_name  VARCHAR(100)",
                "ALTER TABLE car_order        ADD COLUMN IF NOT EXISTS sp_phone VARCHAR(20)",
                "ALTER TABLE contact          ADD COLUMN IF NOT EXISTS sp_code     VARCHAR(20)",
                "ALTER TABLE contact          ADD COLUMN IF NOT EXISTS sp_name     VARCHAR(100)",
                "ALTER TABLE contact          ADD COLUMN IF NOT EXISTS sp_phone    VARCHAR(20)",
                "ALTER TABLE contact          ADD COLUMN IF NOT EXISTS reply_text  TEXT",
                "ALTER TABLE contact          ADD COLUMN IF NOT EXISTS replied_at  TIMESTAMP",
                "ALTER TABLE salesperson      ADD COLUMN IF NOT EXISTS reset_token        VARCHAR(64)",
                "ALTER TABLE salesperson      ADD COLUMN IF NOT EXISTS reset_token_expiry TIMESTAMP",
                "ALTER TABLE loan_application ADD COLUMN IF NOT EXISTS bank_choice        VARCHAR(20) DEFAULT 'republic'",
            ]
            for sql in migrations:
                try:
                    db.session.execute(db.text(sql))
                except Exception:
                    pass
            db.session.commit()

            # Set default 10-year warranty on all solar systems that have none
            try:
                db.session.execute(db.text(
                    "UPDATE solar_system SET warranty = '10 Years' WHERE warranty IS NULL OR warranty = ''"
                ))
                db.session.commit()
            except Exception:
                pass

            if not AdminUser.query.first():
                a = AdminUser(username='admin', full_name='Administrator', email='admin@autopowerdealership.com')
                a.set_password('admin123')
                db.session.add(a)
                db.session.commit()
                print("\n[OK] Default admin created")
                print("   Username : admin")
                print("   Password : admin123")
            print("   WARNING  : Change this password after first login!\n")
        for folder in ('vehicles', 'solar'):
            os.makedirs(os.path.join('static', 'uploads', folder), exist_ok=True)
    except Exception as e:
        print(f"[INIT] DB init warning: {e}")


# Initialise DB on every startup (works with gunicorn / Render too)
try:
    init_db()
except Exception as _e:
    print(f"[INIT] Startup warning: {_e}")

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    print("[START] AutoPower Dealership is running at http://localhost:5000")
    print("[ADMIN] Admin panel at http://localhost:5000/admin")
    print("[SALES] Sales portal at http://localhost:5000/sales/login\n")
    app.run(debug=False, host='0.0.0.0', port=port)
